"""Tests de add_pricing (lib/slides.py). ASCII puro."""

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import lib.slides as s


def test_fmt_miles_agrupa_de_tres_en_tres():
    assert s._fmt_miles(1) == "1"
    assert s._fmt_miles(999) == "999"
    assert s._fmt_miles(1300) == "1.300"
    assert s._fmt_miles(1234567) == "1.234.567"


def test_fmt_eur_entero_sin_decimales():
    assert s._fmt_eur(1300) == "1.300 " + s.EURO
    assert s._fmt_eur(5530) == "5.530 " + s.EURO


def test_fmt_eur_float_sin_parte_decimal_no_pinta_decimales():
    assert s._fmt_eur(1300.0) == "1.300 " + s.EURO


def test_fmt_eur_con_decimales_usa_coma():
    assert s._fmt_eur(1300.5) == "1.300,50 " + s.EURO
    assert s._fmt_eur(0.05) == "0,05 " + s.EURO


def test_fmt_eur_negativo_para_descuentos():
    assert s._fmt_eur(-250) == "-250 " + s.EURO


def test_euro_es_el_simbolo_no_la_cadena_eur():
    assert s.EURO == "\u20AC"


def test_fmt_eur_cero_negativo_no_lleva_signo():
    assert s._fmt_eur(-0.004) == "0 " + s.EURO
    assert s._fmt_eur(-0.0) == "0 " + s.EURO


def test_fmt_eur_acarreo_al_redondear_centimos():
    assert s._fmt_eur(1299.999) == "1.300 " + s.EURO
    assert s._fmt_eur(0.999) == "1 " + s.EURO


def _filas(n):
    return [("Partida %d" % i, 100 * i) for i in range(1, n + 1)]


def test_split_rows_una_pagina_hasta_cinco():
    for n in (1, 3, 5):
        paginas = s._split_rows(_filas(n))
        assert len(paginas) == 1
        assert len(paginas[0]) == n


def test_split_rows_dos_paginas_equilibradas():
    assert [len(p) for p in s._split_rows(_filas(6))] == [3, 3]
    assert [len(p) for p in s._split_rows(_filas(7))] == [4, 3]
    assert [len(p) for p in s._split_rows(_filas(10))] == [5, 5]


def test_split_rows_ninguna_pagina_excede_el_maximo():
    for n in range(1, s.MAX_PRICING_ROWS + 1):
        for pagina in s._split_rows(_filas(n)):
            assert len(pagina) <= s.ROWS_PER_PAGE


def test_split_rows_conserva_el_orden():
    paginas = s._split_rows(_filas(7))
    juntas = paginas[0] + paginas[1]
    assert juntas == _filas(7)


def _prs():
    return s.new_presentation()


def test_add_pricing_devuelve_lista_incluso_con_una_pagina():
    out = s.add_pricing(_prs(), "Inversion", _filas(3))
    assert isinstance(out, list)
    assert len(out) == 1


def test_add_pricing_dos_paginas_con_siete_partidas():
    out = s.add_pricing(_prs(), "Inversion", _filas(7))
    assert len(out) == 2


def test_add_pricing_invoca_el_contador_una_vez_por_pagina():
    llamadas = []

    def contador():
        llamadas.append(len(llamadas) + 1)
        return len(llamadas)

    s.add_pricing(_prs(), "Inversion", _filas(7), page=contador)
    assert llamadas == [1, 2]


def test_add_pricing_rows_vacia_lanza():
    with pytest.raises(ValueError, match="vacia"):
        s.add_pricing(_prs(), "Inversion", [])


def test_add_pricing_mas_de_diez_lanza_en_vez_de_truncar():
    with pytest.raises(ValueError, match="maximo"):
        s.add_pricing(_prs(), "Inversion", _filas(11))


def test_add_pricing_total_numerico_lanza():
    with pytest.raises(TypeError, match="None o str"):
        s.add_pricing(_prs(), "Inversion", _filas(3), total=5530)


def test_add_pricing_page_evaluado_lanza():
    with pytest.raises(TypeError, match="contador"):
        s.add_pricing(_prs(), "Inversion", _filas(3), page=7)


def _textos(slide):
    out = []
    for sp in slide.shapes:
        if sp.has_text_frame:
            out.append(sp.text_frame.text)
    return out


def test_las_filas_pintan_concepto_e_importe_formateado():
    slide = s.add_pricing(_prs(), "Inversion", [("Desarrollo", 2950)])[0]
    textos = _textos(slide)
    assert "Desarrollo" in textos
    assert ("2.950 " + s.EURO) in textos


def test_los_ordinales_continuan_en_la_segunda_pagina():
    paginas = s.add_pricing(_prs(), "Inversion", _filas(7))
    assert "1" in _textos(paginas[0])
    assert "4" in _textos(paginas[0])
    assert "5" in _textos(paginas[1])
    assert "7" in _textos(paginas[1])
    assert "5" not in _textos(paginas[0])


def test_el_total_se_suma_solo():
    slide = s.add_pricing(_prs(), "Inversion",
                          [("A", 1300), ("B", 2950), ("C", 1280)])[0]
    assert ("5.530 " + s.EURO) in _textos(slide)


def test_total_string_se_pinta_tal_cual():
    slide = s.add_pricing(_prs(), "Inversion", _filas(2),
                          total="A convenir")[0]
    assert "A convenir" in _textos(slide)


def test_la_tarjeta_de_total_solo_esta_en_la_ultima_pagina():
    paginas = s.add_pricing(_prs(), "Inversion", _filas(7))
    assert "TOTAL ESTIMADO" not in _textos(paginas[0])
    assert "TOTAL ESTIMADO" in _textos(paginas[1])


def test_la_nota_solo_esta_en_la_ultima_pagina():
    paginas = s.add_pricing(_prs(), "Inversion", _filas(7),
                            note="Precios sin IVA.")
    assert "Precios sin IVA." not in _textos(paginas[0])
    assert "Precios sin IVA." in _textos(paginas[1])


def test_importes_negativos_restan_del_total():
    slide = s.add_pricing(_prs(), "Inversion",
                          [("Web", 1000), ("Descuento", -250)])[0]
    textos = _textos(slide)
    assert ("-250 " + s.EURO) in textos
    assert ("750 " + s.EURO) in textos


def _formas_por_texto(slide, texto):
    return [sp for sp in slide.shapes
            if sp.has_text_frame and sp.text_frame.text == texto]


def _forma_por_texto(slide, texto):
    formas = _formas_por_texto(slide, texto)
    if not formas:
        raise AssertionError("no encuentro la forma con texto %r" % texto)
    return formas[0]


def _forma_en_la_tarjeta(slide, texto):
    """La tarjeta de total vive en la columna derecha. Con una sola partida, el
    total coincide con el importe de esa fila y hay DOS formas con el mismo
    texto: nos quedamos con la de mas a la derecha."""
    formas = _formas_por_texto(slide, texto)
    if not formas:
        raise AssertionError("no encuentro la forma con texto %r" % texto)
    return max(formas, key=lambda sp: sp.left)


def test_las_tres_cajas_de_la_tarjeta_de_total_no_se_solapan():
    # Con 1 o 2 partidas se activa CARD_H_MIN: es el caso en el que la caja de
    # la cifra (1.2in) puede comerse a la etiqueta y a la coletilla.
    for n in (1, 2, 3, 5):
        slide = s.add_pricing(_prs(), "Inversion", _filas(n))[0]
        etiqueta = _forma_en_la_tarjeta(slide, "TOTAL ESTIMADO")
        coletilla = _forma_en_la_tarjeta(slide, "IVA no incluido")
        cifra = _forma_en_la_tarjeta(slide, s._fmt_eur(sum(r[1] for r in _filas(n))))
        # la cifra tiene que estar en la columna de la tarjeta, no en una fila
        assert cifra.left > int(s.MARGIN) + int(s.CONTENT_W) // 2, \
            "n=%d la cifra localizada no esta en la tarjeta" % n
        assert etiqueta.top + etiqueta.height <= cifra.top, "n=%d etiqueta pisa cifra" % n
        assert cifra.top + cifra.height <= coletilla.top, "n=%d cifra pisa coletilla" % n


def test_la_tarjeta_de_total_cabe_en_la_diapositiva():
    for n in (1, 2, 5):
        slide = s.add_pricing(_prs(), "Inversion", _filas(n), note="Precios sin IVA.")[0]
        etiqueta = _forma_en_la_tarjeta(slide, "TOTAL ESTIMADO")
        coletilla = _forma_en_la_tarjeta(slide, "IVA no incluido")
        assert etiqueta.top > 0, "n=%d la tarjeta se sale por arriba" % n
        assert coletilla.top + coletilla.height <= int(s.T.SLIDE_H), \
            "n=%d la tarjeta se sale por abajo" % n


def test_la_nota_no_invade_el_numero_de_pagina():
    from pptx.util import Inches as In
    slide = s.add_pricing(_prs(), "Inversion", _filas(2),
                          note="Nota larga. " * 20)[0]
    nota = _forma_por_texto(slide, "Nota larga. " * 20)
    # _pagenum empieza en SLIDE_W - 1.3in
    assert nota.left + nota.width <= int(s.T.SLIDE_W) - int(In(1.3))
