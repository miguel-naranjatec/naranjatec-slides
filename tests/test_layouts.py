"""Tests de los layouts que validan su entrada. ASCII puro.

El resto de layouts solo dibujan: no hay logica que probar sin abrir PowerPoint.
Aqui solo se fijan los limites que, si se rompen en silencio, hacen perder
contenido (un producto adicional, un punto de la solucion).
"""

import re
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import brand.theme as T
import lib.slides as s


def _prs():
    return s.new_presentation()


def _extra(price=100, image=None):
    return {"price": price, "name": "X", "desc": "y",
            "icon": T.ICON["bolt"], "image": image}


def _punto():
    return ("Titular", "Texto", T.ICON["check"])


def _imgs():
    return [T.img("ecco-working.jpg"), T.img("ecco-office.jpg")]


# --- add_extras -----------------------------------------------------------

def test_add_extras_admite_de_dos_a_cuatro():
    for n in (2, 3, 4):
        slide = s.add_extras(_prs(), "Extras", [_extra()] * n)
        assert slide is not None


def test_add_extras_con_una_sola_lanza():
    with pytest.raises(ValueError, match="entre 2 y 4"):
        s.add_extras(_prs(), "Extras", [_extra()])


def test_add_extras_con_cinco_lanza_en_vez_de_truncar():
    with pytest.raises(ValueError, match="entre 2 y 4"):
        s.add_extras(_prs(), "Extras", [_extra()] * 5)


def test_add_extras_sin_imagen_no_revienta():
    slide = s.add_extras(_prs(), "Extras", [_extra(), _extra()])
    assert slide is not None


def test_add_extras_formatea_el_precio_con_signo_mas():
    slide = s.add_extras(_prs(), "Extras", [_extra(375), _extra(150)])
    textos = [sp.text_frame.text for sp in slide.shapes if sp.has_text_frame]
    assert ("+375 " + s.EURO) in textos
    assert ("+150 " + s.EURO) in textos


# --- add_solution ---------------------------------------------------------

def test_add_solution_admite_de_dos_a_cuatro_puntos():
    for n in (2, 3, 4):
        slide = s.add_solution(_prs(), "Solucion", [_punto()] * n, _imgs())
        assert slide is not None


def test_add_solution_con_cinco_puntos_lanza():
    with pytest.raises(ValueError, match="entre 2 y 4"):
        s.add_solution(_prs(), "Solucion", [_punto()] * 5, _imgs())


def test_add_solution_exige_dos_imagenes():
    with pytest.raises(ValueError, match="2 imagenes"):
        s.add_solution(_prs(), "Solucion", [_punto()] * 2, _imgs()[:1])


def test_add_solution_highlight_anade_una_forma():
    sin = s.add_solution(_prs(), "S", [_punto()] * 3, _imgs())
    con = s.add_solution(_prs(), "S", [_punto()] * 3, _imgs(), highlight=1)
    assert len(con.shapes) == len(sin.shapes) + 1


# --- add_stats_feature ----------------------------------------------------

def _stat(v="+22"):
    return {"value": v, "label": "una etiqueta corta", "icon": T.ICON["shield"]}


def test_add_stats_feature_admite_de_dos_a_cuatro():
    for n in (2, 3, 4):
        slide = s.add_stats_feature(_prs(), "Cifras", [_stat()] * n,
                                    T.img("ecco-working.jpg"))
        assert slide is not None


def test_add_stats_feature_con_cinco_lanza():
    with pytest.raises(ValueError, match="entre 2 y 4"):
        s.add_stats_feature(_prs(), "Cifras", [_stat()] * 5,
                            T.img("ecco-working.jpg"))


def test_ninguna_caja_con_medida_no_positiva():
    # Un alto negativo produce un .pptx que PowerPoint se niega a abrir, y
    # python-pptx no se queja. Este es el caso que lo destapo.
    slide = s.add_stats_feature(_prs(), "Cifras", [_stat()] * 4,
                                T.img("ecco-working.jpg"),
                                subtitle="Con subtitulo, que es el caso justo")
    for sp in slide.shapes:
        assert int(sp.height) > 0
        assert int(sp.width) > 0


def test_check_box_lanza_con_medidas_no_positivas():
    with pytest.raises(ValueError, match="caja invalida"):
        s._check_box(s.Inches(2), s.Inches(-0.5), "prueba")
    with pytest.raises(ValueError, match="caja invalida"):
        s._check_box(0, s.Inches(1), "prueba")


# --- add_next_steps -------------------------------------------------------

def _paso():
    return ("Titular", "Texto corto.", T.ICON["check"])


def test_add_next_steps_admite_de_tres_a_cinco():
    for n in (3, 4, 5):
        slide = s.add_next_steps(_prs(), "Pasos", [_paso()] * n)
        assert slide is not None


def test_add_next_steps_con_dos_lanza():
    with pytest.raises(ValueError, match="entre 3 y 5"):
        s.add_next_steps(_prs(), "Pasos", [_paso()] * 2)


def test_add_next_steps_con_seis_lanza():
    with pytest.raises(ValueError, match="entre 3 y 5"):
        s.add_next_steps(_prs(), "Pasos", [_paso()] * 6)


def test_add_next_steps_pinta_un_arco_y_una_punta_por_hueco():
    # n pasos -> n-1 arcos (formas libres, custGeom) y n-1 puntas (triangulos).
    for n in (3, 4, 5):
        slide = s.add_next_steps(_prs(), "Pasos", [_paso()] * n)
        xml = slide.shapes._spTree.xml
        # OJO: "custGeom" aparece dos veces por forma (apertura y cierre).
        assert xml.count("<a:custGeom>") == n - 1, "arcos con n=%d" % n
        assert xml.count('prst="triangle"') == n - 1, "puntas con n=%d" % n


# --- add_blocks_grid ------------------------------------------------------

def _bloque():
    return ("hero-banner", "Hero", "Imagen, titular y CTA.")


def test_add_blocks_grid_devuelve_siempre_lista():
    out = s.add_blocks_grid(_prs(), "Bloques", [_bloque()] * 3)
    assert isinstance(out, list) and len(out) == 1


def test_add_blocks_grid_pagina_de_doce_en_doce():
    for n, paginas in ((12, 1), (13, 2), (24, 2), (25, 3)):
        out = s.add_blocks_grid(_prs(), "Bloques", [_bloque()] * n)
        assert len(out) == paginas, "n=%d" % n


def test_add_blocks_grid_invoca_el_contador_por_pagina():
    llamadas = []

    def contador():
        llamadas.append(len(llamadas) + 1)
        return len(llamadas)

    s.add_blocks_grid(_prs(), "Bloques", [_bloque()] * 13, page=contador)
    assert llamadas == [1, 2]


def test_add_blocks_grid_vacia_lanza():
    with pytest.raises(ValueError, match="vacia"):
        s.add_blocks_grid(_prs(), "Bloques", [])


def test_add_blocks_grid_page_evaluado_lanza():
    with pytest.raises(TypeError, match="contador"):
        s.add_blocks_grid(_prs(), "Bloques", [_bloque()], page=7)


def test_block_inexistente_lanza_con_pista():
    with pytest.raises(FileNotFoundError, match="make_blocks"):
        T.block("este-bloque-no-existe")


def test_todos_los_svg_tienen_su_png():
    # El SVG es la fuente; el PNG es lo que entra en el .pptx. Si falta uno,
    # el deck revienta al construirse, no al abrirse.
    svgs = {p.stem for p in T.BLOCKS_SRC.glob("*.svg")}
    pngs = {p.stem for p in T.BLOCKS_DIR.glob("*.png")}
    assert svgs, "no hay SVG de bloques"
    assert svgs - pngs == set(), "sin rasterizar: %s" % sorted(svgs - pngs)
    assert pngs - svgs == set(), "PNG huerfanos: %s" % sorted(pngs - svgs)


# --- catalogo de bloques (scripts/gen_blocks.py) ---------------------------

def _gen_blocks():
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))
    import gen_blocks
    return gen_blocks


def test_ningun_slug_repetido_entre_grupos():
    # `BLOCKS` es un dict derivado de `CATALOGO`: si dos grupos declaran el mismo
    # slug, el segundo se come al primero en silencio y perdemos un bloque sin
    # que nadie se entere.
    g = _gen_blocks()
    vistos, repes = set(), []
    for _, grupo in g.CATALOGO:
        for slug in grupo:
            if slug in vistos:
                repes.append(slug)
            vistos.add(slug)
    assert repes == [], "slugs repetidos: %s" % sorted(repes)


def test_todo_bloque_tiene_categoria():
    g = _gen_blocks()
    declarados = sum(len(grupo) for _, grupo in g.CATALOGO)
    assert declarados == len(g.BLOCKS)


def test_cada_svg_generado_corresponde_a_un_bloque_del_catalogo():
    # Un SVG que sobra en disco es un bloque borrado del catalogo cuyo fichero
    # nadie limpio: aparecera en el deck aunque ya no exista la funcion.
    g = _gen_blocks()
    svgs = {p.stem for p in T.BLOCKS_SRC.glob("*.svg")}
    assert svgs == set(g.BLOCKS), "descuadre: %s" % sorted(svgs ^ set(g.BLOCKS))


# --- catalogo de imagenes (brand/imagenes.py) ------------------------------

def _imagenes():
    import brand.imagenes
    return brand.imagenes


def test_ninguna_imagen_se_queda_sin_catalogar():
    # Si alguien mete una foto nueva y no la etiqueta, el agente la elegiria a
    # ciegas por el nombre del fichero. Que es como una propuesta de chocolate
    # acabo con una portada de pulseras.
    IM = _imagenes()
    huerfanas = [r for r, s, _, _ in IM.catalogo() if s == "desconocido"]
    assert huerfanas == [], "sin sector: %s" % huerfanas


def test_el_catalogo_cubre_todas_las_imagenes_del_disco():
    IM = _imagenes()
    en_disco = len(list(IM.IMG_DIR.rglob("*.jpg")))
    assert len(IM.catalogo()) == en_disco


def test_la_advertencia_nunca_calla():
    # Aunque un sector tenga mockups, sigue sin haber foto de producto: la
    # funcion siempre tiene algo que decir, y lo dice.
    IM = _imagenes()
    assert IM.HAY_FOTOGRAFIA_DE_PRODUCTO is False
    for s in IM.sectores():
        aviso = IM.advertencia(s)
        assert aviso and "mockup" in aviso.lower(), s


def test_buscar_filtra_por_sector_y_tipo():
    IM = _imagenes()
    assert "importaco-macbookpro.jpg" in IM.buscar(sector="alimentacion")
    assert IM.buscar(sector="alimentacion", tipo="logo") == []
    assert all("vielong" in r for r in IM.buscar(sector="cosmetica"))
    assert IM.buscar(sector="sector-que-no-existe") == []


# --- diapositivas fijas (content/fijas.py) ---------------------------------

def _fijas():
    import content.fijas
    return content.fijas


def test_cada_fija_construye_su_slide():
    # Una fija rota solo se descubriria generando un deck de cliente.
    F = _fijas()
    for f in F.FIJAS:
        prs = _prs()
        antes = len(prs.slides._sldIdLst)
        f["fn"](prs, page=1, section="Test")
        assert len(prs.slides._sldIdLst) == antes + 1, \
            "la fija %s no anadio una diapositiva" % f["slug"]


def test_el_registro_de_fijas_esta_bien_formado():
    F = _fijas()
    slugs = [f["slug"] for f in F.FIJAS]
    assert len(slugs) == len(set(slugs)), "slugs de fija repetidos: %s" % slugs
    for f in F.FIJAS:
        for clave in ("slug", "nombre", "desc", "acto", "defecto", "fn"):
            assert clave in f, "la fija %s no declara %s" % (f["slug"], clave)
        assert callable(f["fn"])


def test_los_pasos_de_serie_caben_en_el_layout():
    # add_next_steps admite de 3 a 5: anadir un paso obliga a quitar otro.
    F = _fijas()
    assert 3 <= len(F.PASOS) <= 5
    assert len(F.MISION) <= 4, "add_mission pinta 4 features como maximo"


def test_las_fijas_se_pueden_traducir_sin_editar_el_modulo():
    # Un deck en ingles no debe obligar a tocar content/fijas.py.
    F = _fijas()
    prs = _prs()
    F.mision(prs, page=1, title="Our *mission*",
             features=[{"icon": T.ICON["check"], "head": "One partner",
                        "text": "Hosting, development and support."}],
             caption=("We help SMEs go digital", "We get your business."))
    F.proximos_pasos(prs, page=2, title="Next *steps*", subtitle="From sign-off",
                     pasos=[("Sign-off", "You approve.", T.ICON["check"]),
                            ("Design", "We design.", T.ICON["idea"]),
                            ("Launch", "We publish.", T.ICON["bolt"])])
    assert len(prs.slides._sldIdLst) == 2


def test_el_contacto_por_defecto_encaja_en_los_dos_cierres():
    # add_cta espera 2-tuplas y add_thanks 3-tuplas: confundirlos revienta.
    F = _fijas()
    assert all(len(c) == 2 for c in F.contacto_cta())
    assert all(len(c) == 3 for c in F.contacto_thanks())
    assert F.contacto_cta(email="x@y.z")[0][1] == "x@y.z"
    assert F.contacto_cta()[0][1] == F.EMAIL


def test_next_steps_no_pisa_la_descripcion_con_titulares_largos():
    # El bug que el render delato: un titular de dos lineas caia sobre su texto.
    largos = [("Aprobacion del presupuesto por parte del cliente", "d", None),
              ("Aportacion de materiales y contenidos", "d", None),
              ("Lanzamiento", "d", None)]
    largos = [(h, t, T.ICON["check"]) for h, t, _ in largos]
    cortos = [("Diseno", "DESCR", T.ICON["check"])] * 3
    largos = [(h, "DESCR", g) for h, _, g in largos]
    alturas = []
    for pasos in (cortos, largos):
        prs = _prs()
        s.add_next_steps(prs, "T", pasos, page=1)
        cajas = [sh.top for sh in prs.slides[0].shapes
                 if sh.has_text_frame and sh.text_frame.text.strip() == "DESCR"]
        assert len(cajas) == 3, "no encontre las 3 descripciones"
        alturas.append(min(cajas))
    assert alturas[1] > alturas[0], \
        "con titulares de 2 lineas la descripcion no bajo: seguiria solapando"


# --- testimonios (content/testimonios.py) ----------------------------------

def _testimonios():
    import content.testimonios
    return content.testimonios


def test_cada_testimonio_construye_su_slide():
    TT = _testimonios()
    for t in TT.TESTIMONIOS:
        prs = _prs()
        TT.quote(prs, t["slug"], page=1, section="Test")
        assert len(prs.slides._sldIdLst) == 1, t["slug"]


def test_el_extracto_es_literal():
    # La regla que no se puede romper: un extracto solo QUITA texto, nunca lo
    # reescribe. Si alguien "mejora" la frase de un cliente, esto lo caza.
    TT = _testimonios()
    def limpia(s):
        return " ".join(s.split()).strip(" .,;:!?").lower()

    for t in TT.TESTIMONIOS:
        if not t["extracto"]:
            continue
        for trozo in t["extracto"].split("[...]"):
            trozo = limpia(trozo)
            if not trozo:
                continue
            assert trozo in limpia(t["texto"]), \
                "%s: '%s' no aparece literal en la cita" % (t["slug"], trozo[:60])


def test_la_medida_estricta_nunca_cuenta_menos_lineas_que_la_optimista():
    # La raiz de los solapes: _line_count por defecto es OPTIMISTA (+4% de ancho)
    # y quien reserva sitio para un bloque con algo debajo debe medir estricto.
    from pptx.util import Inches, Pt, Emu
    texto = ("Buscabamos una web agil de gestionar que luego podiamos mantener "
             "actualizada de forma independiente y sin depender de nadie.")
    opt = {"size": Pt(23), "italic": True, "font": T.FONT_TITLE_EMPH}
    caja = Emu(int(Inches(7.6)))
    optimista = s._line_count([(texto, opt)], caja)
    estricta = s._line_count([(texto, opt)], caja, tol=s.TOL_ESTRICTA)
    assert s.TOL_ESTRICTA < 1.0
    assert estricta >= optimista, "la medida estricta no puede ser mas confiada"


def test_add_quote_encoge_la_cita_larga_en_vez_de_pisar_al_autor():
    # La cita integra de Importaco (74 palabras) desbordaba la tarjeta y caia
    # sobre el nombre. El cuerpo se elige midiendo, no a ojo.
    TT = _testimonios()
    largo = TT.texto("importaco", corto=False)
    corto = TT.texto("bendorfy")

    def cuerpo_cita(texto):
        prs = _prs()
        s.add_quote(prs, texto, author="A", role="R", page=1)
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and texto[:24] in sh.text_frame.text:
                return sh.text_frame.paragraphs[0].runs[0].font.size.pt
        raise AssertionError("no encontre la cita")

    assert cuerpo_cita(corto) == 23, "una cita corta debe ir al cuerpo mayor"
    assert cuerpo_cita(largo) < 23, "una cita de 74 palabras debe encoger"


def test_por_tema_ordena_por_afinidad():
    TT = _testimonios()
    assert TT.por_tema("hosting") == ["pig-hen", "selva-digital"]
    # El que comparte dos temas va antes que el que comparte uno.
    ranking = TT.por_tema("plazos", "incidencias")
    assert ranking[0] == "importaco"
    assert "valencia-guias" in ranking
    assert TT.por_tema("tema-que-no-existe") == []


def test_el_extracto_solo_existe_si_la_cita_es_larga():
    # Un extracto de una cita que ya cabe holgada es trabajo tirado y una fuente
    # de divergencia entre las dos versiones. El umbral son 25 palabras: por ahi
    # anda la cita de Canon, que a 23pt ocupaba cuatro lineas de la tarjeta.
    TT = _testimonios()
    for t in TT.TESTIMONIOS:
        if t["extracto"]:
            assert len(t["texto"].split()) >= 25, \
                "%s: no necesita extracto" % t["slug"]
            assert len(t["extracto"].split()) < len(t["texto"].split()), \
                "%s: el extracto no acorta nada" % t["slug"]


# --- el skill propuesta-a-deck --------------------------------------------
# El skill le dice a la IA que layouts y que bloques usar. Si uno se renombra o
# se borra, el skill sigue recomendandolo y el deck revienta al construirse.

SKILL = (Path(__file__).resolve().parent.parent / ".claude" / "skills" /
         "propuesta-a-deck" / "SKILL.md")

_SLUG = re.compile(r"`([a-z][a-z0-9]*(?:-[a-z0-9]+)+)`")
_LAYOUT = re.compile(r"`(add_[a-z_]+)`")


def test_el_skill_solo_cita_layouts_que_existen():
    citados = set(_LAYOUT.findall(SKILL.read_text(encoding="utf-8")))
    assert citados, "el skill no cita ningun layout"
    faltan = sorted(nombre for nombre in citados if not hasattr(s, nombre))
    assert faltan == [], "el skill cita layouts inexistentes: %s" % faltan


def test_el_skill_solo_cita_bloques_y_fijas_que_existen():
    texto = SKILL.read_text(encoding="utf-8")
    g = _gen_blocks()
    conocidos = set(g.BLOCKS) | {f["slug"] for f in _fijas().FIJAS}
    # El skill nombra su propio slug en el frontmatter; no es un bloque.
    citados = set(_SLUG.findall(texto)) - {"propuesta-a-deck"}
    assert citados, "el skill no cita ningun bloque"
    faltan = sorted(slug for slug in citados if slug not in conocidos)
    assert faltan == [], "el skill cita slugs inexistentes: %s" % faltan
