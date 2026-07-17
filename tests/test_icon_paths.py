"""Tests de la geometria de los iconos (brand/icon_paths.py).

Se testea el DATO, no el dibujo: si estos invariantes se rompen, el deck sale
mal en silencio y no hay forma de verlo sin abrir PowerPoint.

No requieren fonttools: corren sobre el modulo ya generado.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pytest
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

import brand.icon_paths as IP
import brand.theme as T
import lib.slides as s


def _prs():
    return s.new_presentation()


def _flatten(contour, steps=16):
    """Contorno -> poligono cerrado (las curvas, aplanadas)."""
    pts = [(contour[0][1], contour[0][2])]
    cur = pts[0]
    for seg in contour[1:]:
        k = seg[0]
        if k == "l":
            cur = (seg[1], seg[2])
            pts.append(cur)
        elif k == "q":
            c, e = (seg[1], seg[2]), (seg[3], seg[4])
            for i in range(1, steps + 1):
                t = i / steps
                u = 1 - t
                pts.append((u * u * cur[0] + 2 * u * t * c[0] + t * t * e[0],
                            u * u * cur[1] + 2 * u * t * c[1] + t * t * e[1]))
            cur = e
        elif k == "c":
            c1, c2, e = (seg[1], seg[2]), (seg[3], seg[4]), (seg[5], seg[6])
            for i in range(1, steps + 1):
                t = i / steps
                u = 1 - t
                pts.append((u**3 * cur[0] + 3 * u * u * t * c1[0]
                            + 3 * u * t * t * c2[0] + t**3 * e[0],
                            u**3 * cur[1] + 3 * u * u * t * c1[1]
                            + 3 * u * t * t * c2[1] + t**3 * e[1]))
            cur = e
    if pts[0] != pts[-1]:
        pts.append(pts[0])
    return pts


def _area(poly):
    a = 0.0
    for i in range(len(poly) - 1):
        a += poly[i][0] * poly[i + 1][1] - poly[i + 1][0] * poly[i][1]
    return a / 2.0


def _winding(polys, px, py):
    w = 0
    for poly in polys:
        for i in range(len(poly) - 1):
            x0, y0 = poly[i]
            x1, y1 = poly[i + 1]
            if y0 == y1:
                continue
            if (y0 <= py) and (y1 > py):
                if x0 + (py - y0) / (y1 - y0) * (x1 - x0) > px:
                    w += 1
            elif (y0 > py) and (y1 <= py):
                if x0 + (py - y0) / (y1 - y0) * (x1 - x0) > px:
                    w -= 1
    return w


def test_hay_geometria_para_cada_icono_del_tema():
    # Si theme.ICON gana una clave y nadie regenera los paths, el layout revienta
    # al pintarla. Mejor aqui.
    assert set(IP.ICON_PATHS) == set(T.ICON)


def test_las_coordenadas_son_enteras_y_caben_en_la_caja():
    for name, contours in IP.ICON_PATHS.items():
        for contour in contours:
            for seg in contour:
                for v in seg[1:]:
                    assert isinstance(v, int), "%s: coordenada no entera" % name
                    assert 0 <= v <= IP.ICON_BOX, \
                        "%s: %d fuera de [0, %d]" % (name, v, IP.ICON_BOX)


def test_todos_los_contornos_abren_con_m_y_cierran_con_z():
    for name, contours in IP.ICON_PATHS.items():
        for contour in contours:
            assert contour[0][0] == "m", "%s: contorno sin moveTo" % name
            assert contour[-1] == ("z",), "%s: contorno sin close" % name


def test_no_hay_contornos_degenerados():
    # instantiateVariableFont deja slivers de ancho cero al congelar el eje wght.
    # Con ellos dentro, even-odd y non-zero DEJAN de coincidir (mail era el caso),
    # y entonces el relleno pasa a depender de que regla aplique cada
    # renderizador. El generador los filtra; esto vigila que siga haciendolo.
    for name, contours in IP.ICON_PATHS.items():
        for i, contour in enumerate(contours):
            a = abs(_area(_flatten(contour)))
            assert a >= 50.0, \
                "%s: contorno %d degenerado (area %.1f)" % (name, i, a)


def test_los_iconos_con_agujero_traen_contornos_en_los_dos_sentidos():
    # Un agujero es un contorno con winding OPUESTO al exterior. Si un cambio de
    # fuente los pusiera todos en el mismo sentido, el icono saldria macizo.
    for name in ("globe", "lock", "person", "cert", "mail"):
        signos = {_area(_flatten(c)) > 0 for c in IP.ICON_PATHS[name]}
        assert signos == {True, False}, "%s: no tiene agujeros reales" % name


def test_even_odd_y_non_zero_rellenan_igual():
    # EL invariante que sostiene todo el enfoque. DrawingML no deja elegir regla
    # de relleno, y cada renderizador (PowerPoint Windows, PowerPoint Mac,
    # LibreOffice -que exporta el PDF del cliente-) puede aplicar la suya. Como
    # ninguna zona acumula |winding| > 1, las dos reglas dan el mismo dibujo y da
    # igual cual use cada uno. Si esto se rompe, el icono se vera distinto segun
    # quien lo abra, y es de las cosas que no se detectan sin abrir PowerPoint.
    N = 60
    for name, contours in IP.ICON_PATHS.items():
        polys = [_flatten(c, steps=8) for c in contours]
        for iy in range(N):
            py = (iy + 0.5) * IP.ICON_BOX / N
            for ix in range(N):
                px = (ix + 0.5) * IP.ICON_BOX / N
                w = _winding(polys, px, py)
                assert (w != 0) == (w % 2 != 0), \
                    "%s: even-odd != non-zero en (%.0f, %.0f), winding=%d" \
                    % (name, px, py, w)


def test_icon_shape_pinta_una_forma_con_geometria_propia():
    slide = s._slide(_prs())
    sp = s._icon_shape(slide, Inches(1), Inches(1), Inches(1), T.ICON["lock"])
    xml = sp._element.xml
    assert "<a:custGeom>" in xml
    assert "prstGeom" not in xml          # las dos son excluyentes
    # Un unico <a:path> con TODOS los contornos: es lo que recorta los agujeros.
    assert xml.count("<a:path ") == 1
    assert xml.count("<a:close/>") == len(IP.ICON_PATHS["lock"])
    assert sp.name == "icon:lock"


def test_icon_shape_con_un_glifo_desconocido_revienta():
    slide = s._slide(_prs())
    with pytest.raises(KeyError, match="glifo desconocido"):
        # PUA por codigo, NUNCA como caracter literal (CLAUDE.md). U+E001 existe
        # en la fuente pero no en theme.ICON: sin geometria hay que reventar, no
        # pintar otro icono cualquiera.
        s._icon_shape(slide, Inches(1), Inches(1), Inches(1), chr(0xE001))


def test_el_lado_del_icono_sigue_la_regla_de_26pt_por_pulgada():
    # El icono NO llena su caja: ocupa 26/72 de ella. Venia de cuando era texto y
    # es lo que le da la proporcion dentro de un chip. Clava el no-op visual.
    slide = s._slide(_prs())
    for caja_in in (0.4, 0.62, 0.92):
        sp = s._icon(slide, Inches(1), Inches(1), Inches(caja_in), T.ICON["bolt"])
        esperado = int(round(max(10, round(caja_in * 26)) * 12700))
        assert sp.width == esperado == sp.height, "caja de %.2f in" % caja_in


def test_align_left_pega_el_icono_al_borde_y_center_lo_centra():
    slide = s._slide(_prs())
    caja, x = Inches(0.62), Inches(2)
    izq = s._icon(slide, x, Inches(1), caja, T.ICON["bolt"], align=PP_ALIGN.LEFT)
    cen = s._icon(slide, x, Inches(1), caja, T.ICON["bolt"])
    assert izq.left == int(x)
    assert cen.left == int(x) + (int(caja) - cen.width) // 2


def test_nudge_desplaza_verticalmente_una_fraccion_de_la_caja():
    slide = s._slide(_prs())
    caja, y = Inches(0.62), Inches(1)
    cero = s._icon(slide, Inches(1), y, caja, T.ICON["bolt"], nudge=0.0)
    sube = s._icon(slide, Inches(1), y, caja, T.ICON["bolt"], nudge=-0.04)
    assert cero.top == int(y) + (int(caja) - cero.width) // 2
    assert sube.top == cero.top + int(round(int(caja) * -0.04))


def test_ningun_layout_deja_un_icono_como_texto():
    # El guardian del cambio entero: si alguien vuelve a pintar un glifo con una
    # fuente de iconos, el deck vuelve a depender de que quien lo abre la tenga
    # instalada, y a fallar en PowerPoint for Mac. Ningun run puede llevar un
    # caracter del area de uso privado.
    import content.muestrario as m
    prs = _prs()
    m.build(prs)
    for i, slide in enumerate(prs.slides, 1):
        for sp in slide.shapes:
            if not sp.has_text_frame:
                continue
            for para in sp.text_frame.paragraphs:
                for r in para.runs:
                    malos = [hex(ord(c)) for c in r.text
                             if 0xE000 <= ord(c) <= 0xF8FF]
                    assert not malos, (
                        "diapositiva %d: glifo PUA %s pintado como texto"
                        % (i, malos))
