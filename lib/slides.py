"""Libreria de maquetacion premium con la marca de NaranjaTec.

Estilo "tech premium": fondo blanco, mucho aire, barra superior con logo, tags
tipo pildora, cifras grandes, motivo asterisco/hexagono e imagenes con esquinas
redondeadas. El contenido (textos, rutas de imagen) se pasa como argumentos:
aqui vive solo la maquetacion.

Funciones publicas principales:
    new_presentation()
    add_cover(...)         portada / hero (con imagen opcional a la derecha)
    add_section(...)       separador de seccion (fondo azul)
    add_index(...)         indice de contenidos con tarjetas numeradas
    add_bullets(...)       titulo + vinetas
    add_two_column(...)    dos paneles (reto / solucion)
    add_service_grid(...)  rejilla de tarjetas de servicio
    add_process(...)       pasos numerados (2x2)
    add_image_feature(...) texto + imagen + dato destacado opcional
    add_stats(...)         cifras destacadas
    add_quote(...)         testimonio con avatar
    add_cta(...)           cierre con llamada a la accion
"""

from pathlib import Path

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import brand.theme as T

_BLANK = 6
_LOGO_RATIO = 553 / 120

# Geometria comun
MARGIN = Inches(0.7)
CONTENT_W = Emu(int(T.SLIDE_W) - 2 * int(MARGIN))
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Interlineado exacto de los titulos (factor sobre el tamano de fuente) y huecos
# verticales de la maquetacion aditiva (se suman a la altura real del bloque).
TITLE_LINE_H = 1.2
GAP_AFTER_TITLE = Inches(0.4)   # titulo -> primer bloque de contenido
GAP_PARA = Inches(0.22)         # entre parrafos apilados

# --- Medicion de texto (maquetacion aditiva) -------------------------------- #
# Mapea las familias de theme a un TTF estatico real (en el repo) para medir el
# ancho del texto con PIL y estimar en cuantas lineas se envuelve un parrafo, y
# asi apilar el bloque siguiente con un margen coherente (1 o 2 lineas).
_FONT_DIR = Path(__file__).resolve().parent.parent / "brand" / "assets" / "fonts"
_FORMA = Path(__file__).resolve().parent.parent / "brand" / "assets" / "forma.png"
_FONT_FILE = {
    "Google Sans": "Google_Sans/static/GoogleSans_17pt-Regular.ttf",
    "Google Sans Medium": "Google_Sans/static/GoogleSans_17pt-Medium.ttf",
    "Instrument Sans": "Instrument_Sans/static/InstrumentSans-Regular.ttf",
    "Instrument Sans Medium": "Instrument_Sans/static/InstrumentSans-Medium.ttf",
    "Playfair Display": "Playfair_Display/static/PlayfairDisplay-Italic.ttf",
    "Geist Mono": "Geist_Mono/static/GeistMono-Bold.ttf",
}
_FONT_FILE_BOLD = {"Google Sans": "Google_Sans/static/GoogleSans_17pt-Bold.ttf"}
_MEASURE_PX = 64          # tamano en px al que se cargan las fuentes para medir
_pil_cache = {}


def _pil_font(name, bold):
    """Carga (cacheada) el TTF de la familia a _MEASURE_PX px para medir anchos.
    Devuelve None si no se encuentra (el llamador cae a una estimacion basta)."""
    key = (name, bool(bold))
    if key in _pil_cache:
        return _pil_cache[key]
    rel = (_FONT_FILE_BOLD.get(name) if bold else None) or _FONT_FILE.get(name)
    font = None
    if rel is not None:
        try:
            font = ImageFont.truetype(str(_FONT_DIR / rel), _MEASURE_PX)
        except Exception:
            font = None
    _pil_cache[key] = font
    return font


def _line_count(runs, box_w):
    """Estima en cuantas lineas se envuelve un parrafo (lista de (txt, opt)) en
    una caja de ancho `box_w` (EMU), midiendo palabra a palabra con la fuente
    real de cada run. Todo en PUNTOS: getlength() mide a _MEASURE_PX px y se
    escala a size_pt. Ajuste voraz, como PowerPoint. Minimo 1."""
    # Tolerancia del 4%: PowerPoint suele ajustar algo mas que la medida cruda
    # (variante optica de la fuente, kerning); evita contar de mas una linea.
    box_pt = int(box_w) / 12700.0 * 1.04
    lines, cur = 1, 0.0
    for txt, opt in runs:
        size_pt = int(opt.get("size", T.SIZE_BODY)) / 12700.0
        f = _pil_font(opt.get("font", T.FONT_BODY), opt.get("bold", False))

        def w(s):
            if not s:
                return 0.0
            if f is not None:
                return f.getlength(s) * size_pt / _MEASURE_PX
            return len(s) * size_pt * 0.5           # estimacion basta sin TTF

        for wi, word in enumerate(txt.split(" ")):
            ww = w(word)
            sep = w(" ") if (wi > 0 or cur > 0) else 0.0
            if cur > 0 and cur + sep + ww > box_pt:
                lines += 1
                cur = ww
            else:
                cur += sep + ww
    return max(1, lines)


def _stack_title(slide, runs, x, y, width, size_pt, line_h=TITLE_LINE_H,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Coloca un parrafo de titulo (runs de _emph_runs) con interlineado exacto y
    devuelve la Y inferior REAL segun el numero de lineas medido, para apilar el
    siguiente bloque con un margen coherente. size_pt = tamano base en puntos."""
    lines = _line_count(runs, width)
    h = Emu(int(round(lines * line_h * size_pt * 12700)))
    _text(slide, x, y, width, h, [runs], align=align, anchor=anchor, line_h=line_h)
    return Emu(int(y) + int(h))


# --------------------------------------------------------------------------- #
# Utilidades de bajo nivel
# --------------------------------------------------------------------------- #
def new_presentation():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def _slide(prs, bg=T.BG):
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def _rect(slide, left, top, width, height, fill=None, line=None,
          line_w=Pt(1), shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    sp.shadow.inherit = False
    return sp


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, line_spacing=1.05, space_after=None,
          line_h=None):
    """line_h: si se indica, fija la altura de linea EXACTA en puntos = line_h *
    (mayor tamano del parrafo), en vez del multiplicador `line_spacing` sobre el
    leading natural de la fuente. Da un interlineado predecible y compacto en los
    titulos (p.ej. line_h=1.2) y permite calcular la altura del bloque sin
    depender de las metricas internas de cada fuente."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, str):
            para = [(para, {})]
        if line_h is not None:
            maxpt = max((int(o.get("size", T.SIZE_BODY)) for _t, o in para),
                        default=int(T.SIZE_BODY)) / 12700.0
            p.line_spacing = Pt(line_h * maxpt)
        else:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = space_after
        for txt, opt in para:
            r = p.add_run()
            r.text = txt
            r.font.size = opt.get("size", T.SIZE_BODY)
            name = opt.get("font", T.FONT_BODY)
            # Si la familia ya lleva el peso en el nombre, no aplicamos negrita
            # sintetica encima (evita grosores dobles/inconsistentes).
            weight_named = any(tok in name for tok in T._WEIGHT_TOKENS)
            r.font.bold = False if weight_named else opt.get("bold", False)
            r.font.italic = opt.get("italic", False)
            r.font.name = name
            r.font.color.rgb = opt.get("color", T.GRIS_TEXTO)
            spc = opt.get("spacing")
            if spc:
                r._r.get_or_add_rPr().set("spc", str(int(spc)))
    return box


def _logo(slide, left, top, width, white=False):
    path = T.LOGO_WHITE_PATH if white else T.LOGO_PATH
    # Altura derivada de la proporcion real del PNG (el logo blanco y el de color
    # no comparten ratio), para no deformar ninguno.
    try:
        iw, ih = Image.open(str(path)).size
        ratio = iw / ih
    except Exception:
        ratio = _LOGO_RATIO
    height = Emu(int(int(width) / ratio))
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return height


def _hex_motif(slide, cx, cy, size=Inches(1.15)):
    """Motivo tidy de 3 hexagonos (panal) centrado en (cx, cy). Contenido."""
    s = int(size)
    small = Emu(int(s * 0.82))
    _rect(slide, Emu(int(cx) - s), Emu(int(cy) - int(s * 0.9)), size, size,
          fill=T.AMARILLO, shape=MSO_SHAPE.HEXAGON)
    _rect(slide, Emu(int(cx) + int(s * 0.02)), Emu(int(cy) - int(s * 0.78)),
          small, small, fill=T.AZUL, shape=MSO_SHAPE.HEXAGON)
    _rect(slide, Emu(int(cx) - int(s * 0.5)), Emu(int(cy) + int(s * 0.02)),
          small, small, fill=T.NARANJA, shape=MSO_SHAPE.HEXAGON)


def _img_accent(slide, x, y, w, h, color=T.AMARILLO, off=Inches(0.1), radius=0.05):
    """Bloque de color que asoma sutilmente tras la imagen (marco premium). El
    desplazamiento pequeno hace que marco e imagen se lean como una sola pieza.
    `radius` debe coincidir con el de la imagen para que el filo sea uniforme."""
    _rect(slide, Emu(int(x) + int(off)), Emu(int(y) + int(off)), w, h, fill=color,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def _soft_shadow(shape, alpha=16000, blur=90000, dist=32000):
    """Sombra exterior suave para que las tarjetas 'floten' sobre el fondo."""
    spPr = shape._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="%s"><a:outerShdw blurRad="%d" dist="%d" '
        'dir="5400000" rotWithShape="0"><a:srgbClr val="0B3D66">'
        '<a:alpha val="%d"/></a:srgbClr></a:outerShdw></a:effectLst>'
        % (_A, blur, dist, alpha)))


def _grad_fill(sp, c1, c2, ang_deg=45, a1=None, a2=None):
    """Rellena una forma con un degradado lineal c1 (inicio) -> c2 (fin).
    ang_deg=45 -> de arriba-izquierda a abajo-derecha. a1/a2: opacidad en
    milesimas (None = opaco); permite un relleno tipo cristal (glassmorphism)."""
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)

    def stop(pos, c, a):
        alpha = ('<a:alpha val="%d"/>' % a) if a is not None else ''
        return ('<a:gs pos="%d"><a:srgbClr val="%s">%s</a:srgbClr></a:gs>'
                % (pos, str(c), alpha))
    grad = parse_xml(
        '<a:gradFill xmlns:a="%s" rotWithShape="1"><a:gsLst>%s%s</a:gsLst>'
        '<a:lin ang="%d" scaled="1"/></a:gradFill>'
        % (_A, stop(0, c1, a1), stop(100000, c2, a2), int(ang_deg * 60000)))
    ln = spPr.find(qn("a:ln"))
    ln.addprevious(grad) if ln is not None else spPr.append(grad)


def _edge(sp, color=None, alpha=45000, w=Pt(1)):
    """Borde fino semitransparente (filo de cristal para glassmorphism)."""
    color = color if color is not None else T.BLANCO
    spPr = sp._element.spPr
    old = spPr.find(qn("a:ln"))
    if old is not None:
        spPr.remove(old)
    spPr.append(parse_xml(
        '<a:ln xmlns:a="%s" w="%d"><a:solidFill><a:srgbClr val="%s">'
        '<a:alpha val="%d"/></a:srgbClr></a:solidFill></a:ln>'
        % (_A, int(w), str(color), alpha)))


def _decor(slide, x, y, size, rotation=0, alpha=14000, blur=120000):
    """Decoracion de marca: el isotipo hexagonal grande, girado, muy tenue y
    difuminado, detras del contenido para dar textura. alpha en milesimas
    (14000 = 14% opaco); blur en EMU (0 = sin difuminar)."""
    if not _FORMA.exists():
        return None
    h = Emu(int(int(size) * 1731 / 2000))
    pic = slide.shapes.add_picture(str(_FORMA), x, y, width=size, height=h)
    pic.rotation = rotation
    blip = pic._element.find(".//" + qn("a:blip"))
    if blip is not None:
        blip.append(parse_xml('<a:alphaModFix xmlns:a="%s" amt="%d"/>'
                              % (_A, alpha)))
    if blur:
        pic._element.spPr.append(parse_xml(
            '<a:effectLst xmlns:a="%s"><a:blur rad="%d"/></a:effectLst>'
            % (_A, blur)))
    return pic


def _icon(slide, x, y, size, glyph, color=T.AZUL_OSCURO):
    """Dibuja un icono monolinea (T.FONT_ICON = Material Symbols Outlined 300) centrado
    en la caja, con un pequeno nudge vertical (-4%) para el centrado optico."""
    dy = int(round(int(size) * -0.04))
    b = slide.shapes.add_textbox(x, Emu(int(y) + dy), size, size)
    tf = b.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.name = T.FONT_ICON
    r.font.size = Pt(max(10, round(int(size) / 914400 * 26)))
    r.font.color.rgb = color
    return b


def _icon_chip(slide, x, y, size, glyph, bg=None, fg=T.AMARILLO):
    """Chip de icono premium: cuadrado redondeado (radio contenido) con relleno
    en degradado e icono centrado. bg define el color base del degradado."""
    bg = bg if bg is not None else T.AZUL_OSCURO
    sp = _rect(slide, x, y, size, size, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
               radius=0.18)
    # Degradado solido (opaco) + sombra suave: nitido sobre fondos planos.
    if str(bg) == str(T.AMARILLO):
        _grad_fill(sp, T.AMARILLO_CLARO, T.AMARILLO)
    else:
        _grad_fill(sp, T.AZUL_CHIP, bg)
    _soft_shadow(sp, alpha=30000, blur=45000, dist=22000)
    _icon(slide, x, y, size, glyph, color=fg)


def _pill(slide, x, y, w, h, text, fill=None, text_color=T.AZUL,
          outline=None, size=T.SIZE_SMALL, bold=True):
    sp = _rect(slide, x, y, w, h, fill=fill, line=outline,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.name = T.FONT_MONO
    r.font.color.rgb = text_color
    return sp


def _asterisk(slide, cx, cy, size, color):
    """Motivo de asterisco de 6 puntas (3 barras redondeadas rotadas)."""
    thick = Emu(int(size) // 6)
    for ang in (0, 60, 120):
        bar = _rect(slide, Emu(int(cx) - int(thick) // 2),
                    Emu(int(cy) - int(size) // 2), thick, size, fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        bar.rotation = ang


def _img(slide, path, x, y, w, h, radius=0.045):
    """Inserta una imagen recortada a modo 'cover' con esquinas redondeadas."""
    iw, ih = Image.open(path).size
    target = int(w) / int(h)
    ratio = iw / ih
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if ratio > target:
        crop = (1 - target / ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - ratio / target) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    if radius:
        spPr = pic._element.spPr
        for tag in ("a:prstGeom", "a:custGeom"):
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
        val = int(radius * 100000)
        geom = parse_xml(
            '<a:prstGeom xmlns:a="%s" prst="roundRect"><a:avLst>'
            '<a:gd name="adj" fmla="val %d"/></a:avLst></a:prstGeom>' % (_A, val))
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.addnext(geom)
        else:
            spPr.insert(0, geom)
    return pic


def _gradient_overlay(slide, color=None, a0=94000, a1=46000, ang=0):
    """Rect a pantalla completa con degradado del color de marca: opaco en un
    extremo (donde asienta el texto) y translucido en el otro (deja respirar la
    foto de fondo). alpha en milesimas de %; ang en grados (0 = izq->der)."""
    color = color if color is not None else T.AZUL_OSCURO
    sp = _rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H)  # fill=None -> a:noFill
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    hexv = str(color)
    grad = parse_xml(
        '<a:gradFill xmlns:a="%s" rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '<a:gs pos="52000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '</a:gsLst><a:lin ang="%d" scaled="1"/></a:gradFill>'
        % (_A, hexv, a0, hexv, (a0 + a1) // 2, hexv, a1, int(ang) * 60000))
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return sp


def _topbar(slide, section=""):
    """Barra superior discreta: logo y etiqueta de seccion con tracking."""
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5))
    if section:
        lbl_w = Inches(4.0)
        lx = Emu(int(T.SLIDE_W) - int(MARGIN) - int(lbl_w))
        _text(slide, lx, Inches(0.48), lbl_w, Inches(0.34),
              [[(section.upper(), {"size": Pt(10.5), "color": T.GRIS_SUAVE,
                                   "font": T.FONT_MONO, "spacing": 220})]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _pagenum(slide, page):
    if page is None:
        return
    _text(slide, Emu(int(T.SLIDE_W) - Inches(1.3)),
          Emu(int(T.SLIDE_H) - Inches(0.55)), Inches(0.8), Inches(0.3),
          [[(str(page).zfill(2), {"size": Pt(11), "bold": True,
                                  "color": T.GRIS_SUAVE, "font": T.FONT_MONO})]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _emph_runs(text, size, color=T.AZUL_OSCURO, bold=True,
               emph_color=None, italic=True):
    """Convierte un titulo con marcas *enfasis* en una lista de runs: el texto
    base va en FONT_TITLE (Google Sans, con peso) y lo que quede entre asteriscos
    en FONT_TITLE_EMPH (Playfair Display, cursiva) para un contraste editorial
    sans/serif. Sin asteriscos -> un unico run en Google Sans."""
    emph_color = color if emph_color is None else emph_color
    runs = []
    for i, seg in enumerate(text.split("*")):
        if seg == "":
            continue
        if i % 2:  # segmento entre asteriscos -> enfasis serif
            runs.append((seg, {"size": size, "color": emph_color,
                               "font": T.FONT_TITLE_EMPH, "italic": italic}))
        else:
            runs.append((seg, {"size": size, "color": color,
                               "font": T.FONT_TITLE, "bold": bold}))
    return runs


# --- Presupuestos (add_pricing) ------------------------------------------

EURO = "\u20AC"          # NUNCA el caracter literal: el fichero es ASCII puro
MAX_PRICING_ROWS = 10    # mas partidas -> ValueError (no se truncan en silencio)
ROWS_PER_PAGE = 5        # filas que caben holgadas en una pagina
ROW_H_MAX = Inches(1.05)  # tope de alto de fila: con 1 partida no se estira
CARD_H_MIN = Inches(2.7)  # minimo para que etiqueta, cifra (1.2in) y coletilla no se toquen


def _fmt_miles(entero):
    """1234567 -> '1.234.567' (separador de miles es-ES)."""
    s = str(int(entero))
    grupos = []
    while len(s) > 3:
        grupos.insert(0, s[-3:])
        s = s[:-3]
    grupos.insert(0, s)
    return ".".join(grupos)


def _centimos(valor):
    """Importe -> centimos enteros. Redondear aqui, una sola vez, es lo que
    garantiza que el total nunca contradiga a la suma de las partidas."""
    return int(round(float(valor) * 100))


def _fmt_eur_centimos(centimos):
    """Centimos enteros -> '1.300 EUR' / '1.300,50 EUR' (es-ES)."""
    negativo = centimos < 0
    c = abs(centimos)
    entero, resto = divmod(c, 100)
    if entero == 0 and resto == 0:
        negativo = False
    if resto:
        cuerpo = "%s,%02d" % (_fmt_miles(entero), resto)
    else:
        cuerpo = _fmt_miles(entero)
    return "%s%s %s" % ("-" if negativo else "", cuerpo, EURO)


def _fmt_eur(valor):
    """Importe en es-ES: 1300 -> '1.300 EUR'; 1300.5 -> '1.300,50 EUR'.

    Los enteros no llevan decimales; los decimales van con coma y dos cifras.
    """
    return _fmt_eur_centimos(_centimos(valor))


def _split_rows(rows):
    """Reparte las partidas en paginas de ROWS_PER_PAGE o menos.

    Con 5 o menos, una sola pagina. Con 6-10, dos paginas equilibradas
    (7 -> 4+3), para que ninguna quede casi vacia.
    """
    n = len(rows)
    if n <= ROWS_PER_PAGE:
        return [list(rows)]
    corte = (n + 1) // 2
    return [list(rows[:corte]), list(rows[corte:])]


def _title(slide, title, y=Inches(1.25), width=None, eyebrow=""):
    """Titulo de contenido limpio (sin barra), tipografia como jerarquia.
    Devuelve la Y inferior REAL del titulo (segun 1 o 2 lineas medidas) para que
    el contenido se apile debajo con un margen coherente."""
    width = width or CONTENT_W
    if eyebrow:
        _pill(slide, MARGIN, Emu(int(y) - int(Inches(0.62))), Inches(2.4),
              Inches(0.42), eyebrow, fill=T.AMARILLO_TINT, text_color=T.AZUL_OSCURO,
              size=Pt(11))
    return _stack_title(slide, _emph_runs(title, Pt(32)), MARGIN, y, width, 32)


# --------------------------------------------------------------------------- #
# Diapositivas
# --------------------------------------------------------------------------- #
def add_cover(prs, title, subtitle="", eyebrow="Hosting - Desarrollo - Soporte",
              image=None, cta="www.naranjatec.com", eslogan=T.ESLOGAN):
    slide = _slide(prs)

    # Mismo logo (posicion y tamano) que en el resto del deck.
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5))

    # Imagen a la derecha con marco de color sutil (o motivo hexagonal). El marco
    # va desplazado `off`; se descuenta del tamano para que marco + imagen queden
    # DENTRO de los margenes derecho e inferior y se lean como una sola pieza.
    if image:
        off = Inches(0.1)
        img_x = Inches(7.75)
        img_y = Inches(0.7)
        img_w = Emu(int(T.SLIDE_W) - int(MARGIN) - int(img_x) - int(off))
        img_h = Emu(int(Inches(6.45)) - int(img_y) - int(off))
        _img_accent(slide, img_x, img_y, img_w, img_h, color=T.AMARILLO, off=off,
                    radius=0.025)
        _img(slide, image, img_x, img_y, img_w, img_h, radius=0.025)
        text_w = Inches(6.4)
    else:
        _hex_motif(slide, Inches(10.9), Inches(3.4), Inches(1.7))
        text_w = Inches(9.0)

    # Bloque de texto apilado (flujo aditivo) y centrado en la banda izquierda.
    y = Inches(2.0)
    if eyebrow:
        _pill(slide, MARGIN, y, Inches(2.65), Inches(0.32), eyebrow,
              fill=T.AMARILLO_TINT, text_color=T.AZUL_OSCURO, size=Pt(8.5))
        y = Emu(int(y) + int(Inches(0.66)))
    tb = _stack_title(slide, _emph_runs(title, Pt(36)), MARGIN, y, text_w, 36,
                      line_h=1.14)
    y = Emu(int(tb) + int(Inches(0.28)))
    if subtitle:
        sh = Emu(int(round(_line_count([(subtitle, {"size": Pt(15),
                 "font": T.FONT_BODY})], text_w) * 1.35 * 15 * 12700)))
        _text(slide, MARGIN, y, text_w, sh,
              [[(subtitle, {"size": Pt(15), "color": T.GRIS_SUAVE})]], line_h=1.35)
        y = Emu(int(y) + int(sh) + int(Inches(0.34)))
    if cta:
        _pill(slide, MARGIN, y, Inches(2.55), Inches(0.5), cta,
              fill=T.AZUL, text_color=T.BLANCO, size=Pt(12))

    # Eslogan: linea sutil abajo-izquierda DENTRO de margenes (sin banda a sangre),
    # en Playfair Display cursiva (mismo serif editorial que el enfasis).
    ey = Emu(int(T.SLIDE_H) - int(Inches(0.7)))
    _text(slide, MARGIN, ey, Inches(8.0), Inches(0.4),
          [[(eslogan, {"size": Pt(14), "italic": True, "color": T.GRIS_SUAVE,
                       "font": T.FONT_TITLE_EMPH})]],
          anchor=MSO_ANCHOR.MIDDLE)
    return slide


def add_section(prs, text, number="", image=None):
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    # Foto de fondo a sangre con degradado azul marino encima: da profundidad
    # sin perder legibilidad ni la contencion de color de marca.
    if image:
        _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
        _gradient_overlay(slide)
    # Logo blanco arriba a la izquierda, misma posicion y tamano que el logo
    # superior de las slides de contenido (_topbar): consistencia en todo el deck.
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5), white=True)
    # Motivo hexagonal contenido a la derecha (solo sin foto; con foto estorba)
    if not image:
        _hex_motif(slide, Inches(11.0), Inches(3.75), Inches(1.35))

    # Bloque de titulo anclado ABAJO A LA IZQUIERDA (sin barra vertical). Se mide
    # la altura real del titulo para que 1 o 2 lineas queden sobre la misma base.
    title_w = Inches(9.5)
    runs = _emph_runs(text, T.SIZE_SECTION, color=T.BLANCO, emph_color=T.AMARILLO)
    size_pt = int(T.SIZE_SECTION) / 12700.0
    title_h = Emu(int(round(_line_count(runs, title_w) * TITLE_LINE_H
                            * size_pt * 12700)))
    baseline = Inches(6.55)
    title_y = Emu(int(baseline) - int(title_h))
    kicker = ("SECCION " + number) if number else "SECCION"
    _text(slide, MARGIN, Emu(int(title_y) - int(Inches(0.52))), Inches(6.0),
          Inches(0.45),
          [[(kicker, {"size": Pt(13), "color": T.AMARILLO,
                      "font": T.FONT_MONO, "spacing": 260})]])
    _text(slide, MARGIN, title_y, title_w, title_h, [runs], line_h=TITLE_LINE_H)
    return slide


def add_index(prs, title, items, page=None, highlight=2):
    """Agenda con tarjetas. items: (titulo, desc) o (titulo, desc, glyph). Cada
    tarjeta lleva chip de icono (color) y numero en mono; la tarjeta `highlight`
    se invierte a azul marino como foco. Sin etiqueta superior (redundante)."""
    slide = _slide(prs)
    # Decoracion de marca: isotipo grande, girado y difuminado en la esquina.
    _decor(slide, Inches(9.6), Inches(-2.3), Inches(7.2), rotation=14)
    _topbar(slide)
    tb = _title(slide, title, y=Inches(1.5))

    items = items[:4]
    n = len(items)
    gap = Inches(0.3)
    total_gap = int(gap) * (n - 1)
    card_w = Emu((int(CONTENT_W) - total_gap) // n)
    card_h = Inches(3.15)
    y = Emu(max(int(Inches(2.95)), int(tb) + int(GAP_AFTER_TITLE)))
    pad = Inches(0.34)
    chip = Inches(0.9)
    defaults = [T.ICON["people"], T.ICON["bolt"], T.ICON["cloud"],
                T.ICON["phone"], T.ICON["gear"], T.ICON["check"]]
    for i, it in enumerate(items):
        if isinstance(it, (list, tuple)):
            title_i = it[0]
            desc_i = it[1] if len(it) > 1 else ""
            glyph = it[2] if len(it) > 2 else defaults[i % len(defaults)]
        else:
            title_i, desc_i, glyph = it, "", defaults[i % len(defaults)]
        x = Emu(int(MARGIN) + i * (int(card_w) + int(gap)))
        hot = (i == highlight)
        card = _rect(slide, x, y, card_w, card_h, fill=(T.AZUL_OSCURO if hot
                     else T.BLANCO), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
        _soft_shadow(card, alpha=16000 if hot else 10000)
        ix = Emu(int(x) + int(pad))
        # Chip de icono arriba-izquierda (invertido en la tarjeta destacada).
        if hot:
            _icon_chip(slide, ix, Emu(int(y) + int(pad)), chip, glyph,
                       bg=T.AMARILLO, fg=T.AZUL_OSCURO)
        else:
            _icon_chip(slide, ix, Emu(int(y) + int(pad)), chip, glyph)
        # Numero en mono arriba-derecha (sutil).
        _text(slide, x, Emu(int(y) + int(pad) + int(Inches(0.16))),
              Emu(int(card_w) - int(pad)), Inches(0.5),
              [[("%02d" % (i + 1), {"size": Pt(20), "bold": True,
                 "color": (T.AMARILLO if hot else T.GRIS_BORDE),
                 "font": T.FONT_MONO})]], align=PP_ALIGN.RIGHT)
        # Titulo + descripcion abajo.
        tcol = T.BLANCO if hot else T.AZUL_OSCURO
        dcol = T.AZUL_TINT if hot else T.GRIS_SUAVE
        _text(slide, ix, Emu(int(y) + int(card_h) - int(Inches(1.18))),
              Emu(int(card_w) - 2 * int(pad)), Inches(0.55),
              [[(title_i, {"size": Pt(16), "color": tcol, "font": T.FONT_HEAD})]])
        if desc_i:
            _text(slide, ix, Emu(int(y) + int(card_h) - int(Inches(0.6))),
                  Emu(int(card_w) - 2 * int(pad)), Inches(0.5),
                  [[(desc_i, {"size": Pt(11), "color": dcol, "font": T.FONT_BODY})]])
    _pagenum(slide, page)
    return slide


def add_bullets(prs, title, bullets, image=None, subtitle="", page=None,
                section=""):
    """Dos zonas: titulo (+ subtitulo opcional en serif) y lista (chip de icono
    + titulo + detalle) a la izquierda e imagen a sangre a la derecha. bullets:
    (head, detail) o (head, detail, glyph). Sin imagen usa todo el ancho."""
    slide = _slide(prs)
    if image:
        ix = Inches(7.7)
        _img(slide, image, ix, 0, Emu(int(T.SLIDE_W) - int(ix)), T.SLIDE_H,
             radius=0)
        _topbar(slide)  # solo logo (la etiqueta chocaria con la imagen)
        text_w = Emu(int(ix) - int(MARGIN) - int(Inches(0.55)))
    else:
        _topbar(slide, section)
        text_w = CONTENT_W
    tb = _stack_title(slide, _emph_runs(title, Pt(30)), MARGIN, Inches(1.5),
                      text_w, 30)
    content_top = tb
    if subtitle:
        sy = Emu(int(tb) + int(Inches(0.12)))
        _text(slide, MARGIN, sy, text_w, Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        content_top = Emu(int(sy) + int(Inches(0.5)))

    items = bullets[:4]
    defaults = [T.ICON["check"], T.ICON["bolt"], T.ICON["shield"], T.ICON["gear"]]
    y0 = Emu(max(int(Inches(3.0)), int(content_top) + int(GAP_AFTER_TITLE)))
    step = Emu((int(Inches(6.7)) - int(y0)) // max(len(items), 1))
    chip = Inches(0.66)
    for i, item in enumerate(items):
        if isinstance(item, (list, tuple)):
            head = item[0]
            detail = item[1] if len(item) > 1 else ""
            glyph = item[2] if len(item) > 2 else defaults[i % len(defaults)]
        else:
            head, detail, glyph = item, "", defaults[i % len(defaults)]
        iy = Emu(int(y0) + i * int(step))
        _icon_chip(slide, MARGIN, iy, chip, glyph)
        hx = Emu(int(MARGIN) + int(chip) + int(Inches(0.35)))
        hw = Emu(int(MARGIN) + int(text_w) - int(hx))
        # Titulo + detalle como un unico bloque centrado con el chip: pareja
        # compacta (space_after pequeno) en vez de un hueco fijo suelto.
        runs = [[(head, {"size": Pt(16), "color": T.AZUL_OSCURO,
                         "font": T.FONT_HEAD})]]
        if detail:
            runs.append([(detail, {"size": Pt(12), "color": T.GRIS_SUAVE,
                                   "font": T.FONT_BODY})])
        _text(slide, hx, iy, hw, chip, runs, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.12, space_after=Pt(3))
    _white_pagenum(slide, page) if image else _pagenum(slide, page)
    return slide


def add_two_column(prs, title, left, right, subtitle="", page=None, section=""):
    """Slide dividido en dos mitades a sangre (reto -> solucion): izquierda clara,
    derecha azul marino. Cada mitad: chip de icono + heading + lista con checks.
    left/right: {"heading", "items", "icon" (opcional)}. subtitle: serif opcional."""
    slide = _slide(prs)
    mid = Emu(int(T.SLIDE_W) // 2)
    # Dos mitades a sangre a TODO el alto que parten la diapositiva.
    _rect(slide, 0, 0, mid, T.SLIDE_H, fill=T.GRIS_FONDO)
    _rect(slide, mid, 0, Emu(int(T.SLIDE_W) - int(mid)), T.SLIDE_H,
          fill=T.AZUL_OSCURO)
    # Logo sobre la mitad clara + etiqueta en blanco sobre la mitad oscura.
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5))
    if section:
        lbl_w = Inches(4.0)
        _text(slide, Emu(int(T.SLIDE_W) - int(MARGIN) - int(lbl_w)), Inches(0.48),
              lbl_w, Inches(0.34),
              [[(section.upper(), {"size": Pt(10.5), "color": T.BLANCO,
                                   "font": T.FONT_MONO, "spacing": 220})]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # Titulo + subtitulo dentro de la mitad izquierda.
    left_w = Emu(int(mid) - int(MARGIN) - int(Inches(0.65)))
    tb = _stack_title(slide, _emph_runs(title, Pt(28)), MARGIN, Inches(1.35),
                      left_w, 28)
    if subtitle:
        _text(slide, MARGIN, Emu(int(tb) + int(Inches(0.12))), left_w, Inches(0.5),
              [[(subtitle, {"size": Pt(13), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])

    # Contenido de cada mitad centrado en la banda inferior (misma base para las
    # dos, para que los encabezados queden enfrentados).
    band_top = int(Inches(2.85))
    body_h = int(T.SLIDE_H) - band_top
    chip = Inches(0.72)
    item_h = Inches(0.62)
    head_gap = Inches(1.15)

    def half(cx, cw, data, dark, default_icon):
        items = data["items"][:4]
        total = int(head_gap) + len(items) * int(item_h)
        gy = Emu(band_top + (body_h - total) // 2)
        if dark:
            _icon_chip(slide, cx, gy, chip, data.get("icon", default_icon),
                       bg=T.AMARILLO, fg=T.AZUL_OSCURO)
        else:
            _icon_chip(slide, cx, gy, chip, data.get("icon", default_icon))
        hx = Emu(int(cx) + int(chip) + int(Inches(0.3)))
        _text(slide, hx, gy, Emu(int(cw) - int(chip) - int(Inches(0.3))), chip,
              [[(data["heading"], {"size": Pt(18),
                 "color": (T.BLANCO if dark else T.AZUL_OSCURO),
                 "font": T.FONT_HEAD})]], anchor=MSO_ANCHOR.MIDDLE)
        iy = Emu(int(gy) + int(head_gap))
        m = Inches(0.58)
        for it in items:
            _icon(slide, cx, Emu(int(iy) + (int(item_h) - int(m)) // 2), m,
                  T.ICON["tick"], color=(T.AMARILLO if dark else T.AZUL))
            _text(slide, Emu(int(cx) + int(Inches(0.7))), iy,
                  Emu(int(cw) - int(Inches(0.7))), item_h,
                  [[(it, {"size": Pt(14),
                          "color": (T.BLANCO if dark else T.GRIS_TEXTO),
                          "font": T.FONT_BODY})]], anchor=MSO_ANCHOR.MIDDLE)
            iy = Emu(int(iy) + int(item_h))

    pad = Inches(0.65)
    half(MARGIN, Emu(int(mid) - int(MARGIN) - int(pad)), left, False,
         T.ICON["pulse"])
    half(Emu(int(mid) + int(pad)),
         Emu(int(T.SLIDE_W) - int(mid) - int(pad) - int(MARGIN)), right, True,
         T.ICON["shield"])
    _white_pagenum(slide, page)
    return slide


def add_service_grid(prs, title, items, subtitle="", page=None, section=""):
    """Rejilla compacta de servicios (sin cajas): chip de icono + titulo +
    descripcion, en 3 columnas y hasta 3 filas (9 items). items: {title, desc,
    icon}. subtitle: serif opcional."""
    slide = _slide(prs)
    # Varias formas de marca en el fondo (esquinas), giradas, tenues y difuminadas.
    _decor(slide, Inches(10.2), Inches(-2.9), Inches(6.4), rotation=16, alpha=9000)
    _decor(slide, Inches(-2.9), Inches(4.6), Inches(6.0), rotation=-22, alpha=8000)
    _decor(slide, Inches(11.6), Inches(5.6), Inches(4.6), rotation=40, alpha=7000)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.4))
    content_top = tb
    if subtitle:
        sy = Emu(int(tb) + int(Inches(0.12)))
        _text(slide, MARGIN, sy, Inches(9.0), Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        content_top = Emu(int(sy) + int(Inches(0.5)))

    items = items[:9]
    n = len(items)
    cols = 3
    rows = (n + cols - 1) // cols
    gap_x = Inches(0.55)
    grid_top = max(int(Inches(2.75)), int(content_top) + int(Inches(0.35)))
    grid_bottom = int(Inches(7.0))
    col_w = Emu((int(CONTENT_W) - (cols - 1) * int(gap_x)) // cols)
    row_h = (grid_bottom - grid_top) // max(rows, 1)
    chip = Inches(0.46)
    gap = Inches(0.24)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = Emu(int(MARGIN) + c * (int(col_w) + int(gap_x)))
        y = Emu(grid_top + r * row_h)
        _icon_chip(slide, x, y, chip, it.get("icon", T.ICON["check"]))
        # Titulo y descripcion alineados al mismo margen (a la derecha del chip).
        hx = Emu(int(x) + int(chip) + int(gap))
        hw = Emu(int(x) + int(col_w) - int(hx))
        _text(slide, hx, y, hw, chip,
              [[(it["title"], {"size": Pt(14.5), "color": T.AZUL_OSCURO,
                               "font": T.FONT_HEAD})]], anchor=MSO_ANCHOR.MIDDLE)
        if it.get("desc"):
            # Descripcion justo bajo el titulo (no bajo el chip): gap minimo.
            dy = Emu(int(y) + int(chip) - int(Inches(0.11)))
            _text(slide, hx, dy, hw, Emu(int(grid_top) + (r + 1) * row_h - int(dy)),
                  [[(it["desc"], {"size": Pt(10.5), "color": T.GRIS_SUAVE,
                                  "font": T.FONT_BODY})]], line_spacing=1.2)
    _pagenum(slide, page)
    return slide


def add_process(prs, title, steps, page=None, section=""):
    """Pasos numerados en rejilla 2x2. steps: lista de (titulo, desc)."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.5))

    cols = 2
    card_w = Inches(5.8)
    card_h = Inches(1.75)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)
    start_y = Emu(max(int(Inches(2.7)), int(tb) + int(GAP_AFTER_TITLE)))
    for i, st in enumerate(steps[:4]):
        title_i, desc_i = (st if isinstance(st, (list, tuple)) else (st, ""))
        row, col = divmod(i, cols)
        x = Emu(int(MARGIN) + col * (int(card_w) + int(gap_x)))
        y = Emu(int(start_y) + row * (int(card_h) + int(gap_y)))
        card = _rect(slide, x, y, card_w, card_h, fill=T.GRIS_FONDO,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        _soft_shadow(card, alpha=9000)
        _rect(slide, Emu(int(x) + int(Inches(0.4))),
              Emu(int(y) + int(Inches(0.42))), Inches(0.9), Inches(0.9),
              fill=T.AZUL_OSCURO, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        _text(slide, Emu(int(x) + int(Inches(0.4))),
              Emu(int(y) + int(Inches(0.42))), Inches(0.9), Inches(0.9),
              [[("%02d" % (i + 1), {"size": Pt(20), "bold": True,
                                    "color": T.AMARILLO, "font": T.FONT_NUM})]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(int(x) + int(Inches(1.55))),
              Emu(int(y) + int(Inches(0.4))),
              Emu(int(card_w) - int(Inches(1.9))), Inches(0.5),
              [[(title_i, {"size": Pt(16), "bold": True,
                           "color": T.AZUL_OSCURO, "font": T.FONT_HEAD})]])
        _text(slide, Emu(int(x) + int(Inches(1.55))),
              Emu(int(y) + int(Inches(0.9))),
              Emu(int(card_w) - int(Inches(1.9))), Inches(0.7),
              [[(desc_i, {"size": Pt(12), "color": T.GRIS_SUAVE})]])
    _pagenum(slide, page)
    return slide


def add_image_feature(prs, title, body, image, stat=None, side="right",
                      subtitle="", page=None, section=""):
    """Texto + imagen. body: str o lista de parrafos. stat: (valor, etiqueta).
    subtitle: serif opcional bajo el titulo."""
    slide = _slide(prs)
    _topbar(slide, section)

    img_w = Inches(5.4)
    img_h = Inches(4.7)
    img_y = Inches(1.7)
    if side == "right":
        img_x = Emu(int(T.SLIDE_W) - int(MARGIN) - int(img_w))
        text_x = MARGIN
    else:
        img_x = MARGIN
        text_x = Emu(int(MARGIN) + int(img_w) + int(Inches(0.6)))
    text_w = Emu(int(T.SLIDE_W) - int(MARGIN) - int(img_w) - int(MARGIN) - int(Inches(0.6)))

    # Marco de color sutil como en la portada (filo pegado a la imagen).
    _img_accent(slide, img_x, img_y, img_w, img_h, off=Inches(0.1), radius=0.025)
    _img(slide, image, img_x, img_y, img_w, img_h, radius=0.025)

    # Titulo medido + cuerpo apilado debajo con margen coherente (1 o 2 lineas
    # de titulo dejan el mismo hueco al cuerpo; sin barra bajo el titulo).
    tb = _stack_title(slide, _emph_runs(title, Pt(30)), text_x, Inches(2.4),
                      text_w, 30)
    body_y = Emu(int(tb) + int(GAP_AFTER_TITLE))
    if subtitle:
        _text(slide, text_x, Emu(int(tb) + int(Inches(0.12))), text_w, Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        body_y = Emu(int(tb) + int(Inches(0.68)))
    if isinstance(body, str):
        body = [body]
    paras = [[(b, {"size": Pt(14), "color": T.GRIS_SUAVE, "font": T.FONT_BODY})]
             for b in body]
    _text(slide, text_x, body_y, text_w, Inches(2.6),
          paras, line_spacing=1.3, space_after=Pt(8))

    if stat:
        val, label = stat
        sc_w = Inches(2.5)
        sc_h = Inches(1.3)
        sc_x = Emu(int(img_x) - int(Inches(0.6)))
        sc_y = Emu(int(img_y) + int(img_h) - int(sc_h) - int(Inches(0.3)))
        _rect(slide, sc_x, sc_y, sc_w, sc_h, fill=T.AMARILLO,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        _text(slide, sc_x, Emu(int(sc_y) + int(Inches(0.2))), sc_w, Inches(0.7),
              [[(val, {"size": Pt(34), "bold": True, "color": T.AZUL_OSCURO,
                       "font": T.FONT_NUM})]], align=PP_ALIGN.CENTER)
        _text(slide, sc_x, Emu(int(sc_y) + int(Inches(0.85))), sc_w, Inches(0.4),
              [[(label, {"size": Pt(11), "bold": True, "color": T.AZUL_OSCURO})]],
              align=PP_ALIGN.CENTER)
    _pagenum(slide, page)
    return slide


def add_stats(prs, stats, title="", subtitle="", page=None, section=""):
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.5)) if title else None
    if subtitle and tb is not None:
        _text(slide, MARGIN, Emu(int(tb) + int(Inches(0.12))), Inches(9.0),
              Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])

    # Tarjetas que flotan: numero grande + acento dorado + etiqueta, centrado.
    content_top = int(tb) if tb is not None else int(Inches(2.2))
    n = max(len(stats), 1)
    gap = int(Inches(0.4))
    card_w = Emu((int(CONTENT_W) - (n - 1) * gap) // n)
    card_h = int(Inches(2.8))
    band_top = max(content_top + int(Inches(0.55)), int(Inches(3.0)))
    band_bot = int(T.SLIDE_H) - int(Inches(0.7))
    cy = band_top + ((band_bot - band_top) - card_h) // 2
    for i, s in enumerate(stats):
        x = Emu(int(MARGIN) + i * (int(card_w) + gap))
        card = _rect(slide, x, Emu(cy), card_w, Emu(card_h), fill=T.BLANCO,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        _soft_shadow(card, alpha=11000)
        _text(slide, x, Emu(cy + int(Inches(0.5))), card_w, Inches(1.15),
              [[(s["value"], {"size": Pt(52), "bold": True,
                              "color": T.AZUL_OSCURO, "font": T.FONT_NUM})]],
              align=PP_ALIGN.CENTER)
        aw = int(Inches(0.7))
        _rect(slide, Emu(int(x) + (int(card_w) - aw) // 2),
              Emu(cy + int(Inches(1.62))), Emu(aw), Pt(3.5), fill=T.AMARILLO)
        _text(slide, Emu(int(x) + int(Inches(0.3))), Emu(cy + int(Inches(1.85))),
              Emu(int(card_w) - int(Inches(0.6))), Inches(0.7),
              [[(s["label"], {"size": Pt(13.5), "color": T.GRIS_SUAVE,
                              "font": T.FONT_BODY})]], align=PP_ALIGN.CENTER,
              line_spacing=1.2)
    _pagenum(slide, page)
    return slide


def add_quote(prs, quote, author, role="", avatar=None, stars=5, image=None,
              page=None, section=""):
    """Testimonio sobre foto a sangre con velo azul marino y tarjeta central
    glassmorfista (translucida, filo de cristal): valoracion en estrellas
    doradas, cita en serif italica y autor (avatar + nombre + cargo), en blanco.
    El cristal se ve porque hay foto detras (sin foto queda azul marino solido)."""
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    if image:
        _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
        _tint(slide, 0, 0, T.SLIDE_W, T.SLIDE_H, color=T.AZUL_OSCURO,
              alpha=46000)
    # Barra superior en version blanca (sobre foto).
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5), white=True)
    if section:
        lbl_w = Inches(4.0)
        lx = Emu(int(T.SLIDE_W) - int(MARGIN) - int(lbl_w))
        _text(slide, lx, Inches(0.48), lbl_w, Inches(0.34),
              [[(section.upper(), {"size": Pt(10.5), "color": T.GRIS_BORDE,
                                   "font": T.FONT_MONO, "spacing": 220})]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Tarjeta glassmorfista centrada: copia del fondo recortada y desenfocada
    # (backdrop-blur) + velo navy translucido + filo de cristal muy tenue.
    cw, chh = Inches(9.4), Inches(4.6)
    cx = Emu((int(T.SLIDE_W) - int(cw)) // 2)
    cyt = int(Inches(1.5))
    if image:
        _glass_card(slide, image, cx, Emu(cyt), cw, Emu(int(chh)), radius=0.05,
                    blur=95000)
    card = _tint(slide, cx, Emu(cyt), cw, Emu(int(chh)), color=T.AZUL_OSCURO,
                 alpha=38000, radius=0.05)
    _edge(card, T.BLANCO, alpha=15000, w=Pt(0.75))
    _soft_shadow(card, alpha=34000, blur=70000, dist=22000)
    pad = int(Inches(0.9))
    inx = Emu(int(cx) + pad)
    inw = Emu(int(cw) - 2 * pad)
    # Estrellas de valoracion: formas nativas rellenas (doradas, centradas).
    # Se usan shapes (no glifo) para que salgan RELLENAS y sin depender de fuente.
    if stars:
        nst = int(stars)
        ssz = int(Inches(0.26))
        sgap = int(Inches(0.12))
        total = nst * ssz + (nst - 1) * sgap
        sx = int(cx) + (int(cw) - total) // 2
        sy = cyt + int(Inches(0.62))
        for i in range(nst):
            _rect(slide, Emu(sx + i * (ssz + sgap)), Emu(sy), Emu(ssz),
                  Emu(ssz), fill=T.AMARILLO, shape=MSO_SHAPE.STAR_5_POINT)
    # Cita en serif italica, centrada, en blanco.
    _text(slide, inx, Emu(cyt + int(Inches(1.25))), inw, Inches(1.5),
          [[('"' + quote + '"', {"size": Pt(23), "italic": True,
                                 "color": T.BLANCO,
                                 "font": T.FONT_TITLE_EMPH})]],
          align=PP_ALIGN.CENTER, line_h=1.3)
    # Autor: avatar centrado + nombre + cargo.
    av = int(Inches(0.82))
    ay = cyt + int(Inches(2.7))
    if avatar:
        _img(slide, avatar, Emu((int(T.SLIDE_W) - av) // 2), Emu(ay), Emu(av),
             Emu(av), radius=0.5)
    ty = ay + av + int(Inches(0.16))
    _text(slide, inx, Emu(ty), inw, Inches(0.4),
          [[(author, {"size": Pt(15), "bold": True, "color": T.BLANCO,
                      "font": T.FONT_HEAD})]], align=PP_ALIGN.CENTER)
    if role:
        _text(slide, inx, Emu(ty + int(Inches(0.4))), inw, Inches(0.4),
              [[(role, {"size": Pt(12), "color": T.GRIS_BORDE})]],
              align=PP_ALIGN.CENTER)
    _white_pagenum(slide, page)
    return slide


def add_cta(prs, headline, contact, subtext="", image=None):
    """Cierre / llamada a la accion. Con `image`: foto a sangre + degradado azul
    marino (calidad de separador); sin ella, azul marino solido con hexagonos.
    contact: lista de (icono, valor) -> se pintan con icono dorado."""
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    if image:
        _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
        _gradient_overlay(slide)
    else:
        _hex_motif(slide, Inches(11.0), Inches(4.3), Inches(1.65))
    _logo(slide, MARGIN, Inches(0.42), Inches(1.5), white=True)

    # Bloque alineado al margen (sin barra ni acento).
    y = int(Inches(2.35))
    hb = _stack_title(slide, _emph_runs(headline, Pt(40), color=T.BLANCO,
                                        emph_color=T.AMARILLO), MARGIN, Emu(y),
                      Inches(9.0), 40)
    y = int(hb) + int(Inches(0.3))
    if subtext:
        _text(slide, MARGIN, Emu(y), Inches(8.5), Inches(0.6),
              [[(subtext, {"size": T.SIZE_SUBTITLE, "color": T.AMARILLO})]],
              line_h=1.2)
        y += int(Inches(0.9))
    else:
        y += int(Inches(0.4))
    # Contactos: icono dorado + valor en blanco (peso sutil, Instrument Sans).
    row_h = int(Inches(0.62))
    isz = int(Inches(0.5))
    for icon, value in contact:
        _icon(slide, MARGIN, Emu(y + (row_h - isz) // 2), Emu(isz), icon,
              color=T.AMARILLO)
        _text(slide, Emu(int(MARGIN) + int(Inches(0.8))), Emu(y), Inches(7.0),
              Emu(row_h), [[(value, {"size": Pt(15), "color": T.BLANCO,
                                     "font": T.FONT_BODY})]],
              anchor=MSO_ANCHOR.MIDDLE)
        y += row_h
    return slide


# --------------------------------------------------------------------------- #
# Layouts adicionales: declaracion, linea de tiempo, comparativa, galerias
# --------------------------------------------------------------------------- #
def _scrim(slide, color=None, stops=None, ang=5400000):
    """Velo degradado a pantalla completa (por defecto oscuro arriba ->
    transparente) para asentar texto blanco sobre una foto. ang=5400000 -> vertical."""
    color = color if color is not None else T.AZUL_OSCURO
    sp = _rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H)
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    hexv = str(color)
    stops = stops or [(0, 86000), (40000, 8000), (100000, 0)]
    gs = "".join('<a:gs pos="%d"><a:srgbClr val="%s"><a:alpha val="%d"/>'
                 '</a:srgbClr></a:gs>' % (p, hexv, a) for p, a in stops)
    grad = parse_xml('<a:gradFill xmlns:a="%s" rotWithShape="1"><a:gsLst>%s'
                     '</a:gsLst><a:lin ang="%d" scaled="1"/></a:gradFill>'
                     % (_A, gs, ang))
    ln = spPr.find(qn("a:ln"))
    ln.addprevious(grad) if ln is not None else spPr.append(grad)
    return sp


def _white_pagenum(slide, page):
    if page is None:
        return
    _text(slide, Emu(int(T.SLIDE_W) - int(Inches(1.3))),
          Emu(int(T.SLIDE_H) - int(Inches(0.62))), Inches(0.8), Inches(0.3),
          [[(str(page).zfill(2), {"size": Pt(11), "bold": True,
             "color": T.BLANCO, "font": T.FONT_MONO})]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_statement(prs, text, image, author="", page=None):
    """Declaracion / cita a sangre completa: foto de fondo, velo oscuro arriba,
    icono de bocadillo dorado y la frase en blanco (admite *enfasis*)."""
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
    _scrim(slide, stops=[(0, 88000), (42000, 10000), (100000, 0)])
    _icon(slide, MARGIN, Inches(0.86), Inches(0.92), T.ICON["quote"],
          color=T.AMARILLO)
    tx = Emu(int(MARGIN) + int(Inches(1.2)))
    tw = Emu(int(T.SLIDE_W) - int(tx) - int(MARGIN))
    bottom = _stack_title(
        slide, _emph_runs(text, Pt(25), color=T.BLANCO, emph_color=T.AMARILLO),
        tx, Inches(0.92), tw, 25, line_h=1.3)
    if author:
        _text(slide, tx, Emu(int(bottom) + int(Inches(0.16))), tw, Inches(0.4),
              [[("/ " + author, {"size": Pt(12.5), "color": T.AMARILLO,
                                 "font": T.FONT_MONO, "spacing": 40})]])
    _logo(slide, MARGIN, Emu(int(T.SLIDE_H) - int(Inches(0.95))), Inches(1.4),
          white=True)
    _white_pagenum(slide, page)
    return slide


def _dashed_vline(slide, x, y1, y2, color=None, w=Pt(1.5)):
    """Eje vertical punteado de la linea de tiempo (guion via XML: python-pptx
    1.x no expone el enum de dash)."""
    color = color if color is not None else T.GRIS_BORDE
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    cn.line.color.rgb = color
    cn.line.width = w
    lnEl = cn.line._get_or_add_ln()
    lnEl.append(parse_xml('<a:prstDash xmlns:a="%s" val="sysDot"/>' % _A))
    cn.shadow.inherit = False
    return cn


def add_timeline(prs, title, milestones, subtitle="", section=""):
    """Linea de tiempo VERTICAL MULTIPAGINA: genera una diapositiva por hito,
    revelando los nodos del eje progresivamente. En cada pagina el hito actual se
    marca con un aro dorado; los anteriores quedan como contexto y los siguientes
    como puntos tenues. milestones: [{year, label, text, icon}] (alternan lados).
    Devuelve la lista de diapositivas creadas."""
    ms = milestones
    n = max(len(ms), 1)
    accents = [T.AZUL, T.AZUL_OSCURO, T.NARANJA]
    cx = int(T.SLIDE_W) // 2
    node_cy = int(T.SLIDE_H) // 2       # nodo actual centrado en la pagina
    node = int(Inches(1.2))
    # Formas de marca de fondo, en posiciones distintas por pagina (cada set es
    # (x, y, tamano, rotacion)); refuerzan la sensacion de recorrer el timeline.
    decor_sets = [
        [(Inches(-2.7), Inches(-2.4), Inches(6.8), 16),
         (Inches(10.6), Inches(4.4), Inches(5.4), -22)],
        [(Inches(9.8), Inches(-3.1), Inches(7.0), 26),
         (Inches(-2.5), Inches(4.0), Inches(5.6), -12)],
        [(Inches(-1.9), Inches(4.6), Inches(6.0), 40),
         (Inches(10.9), Inches(-2.6), Inches(5.2), 8)],
    ]
    created = []
    for i, m in enumerate(ms):
        slide = _slide(prs)
        for dx, dy, dsz, drot in decor_sets[i % len(decor_sets)]:
            _decor(slide, dx, dy, dsz, rotation=drot, alpha=8000)
        _topbar(slide, section)
        _title(slide, title, y=Inches(1.35))
        if subtitle:
            _text(slide, MARGIN, Inches(2.0), Inches(9.0), Inches(0.5),
                  [[(subtitle, {"size": Pt(14), "italic": True,
                                "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])

        # Eje punteado como un viewport: se prolonga hacia arriba si hay pasado y
        # hacia abajo si queda futuro (a sangre por el borde correspondiente).
        y_top = 0 if i > 0 else node_cy
        y_bot = int(T.SLIDE_H) if i < n - 1 else node_cy
        _dashed_vline(slide, Emu(cx), Emu(y_top), Emu(y_bot))

        # Nodo actual (centrado) con icono.
        accent = accents[i % len(accents)]
        nx, ny = cx - node // 2, node_cy - node // 2
        _rect(slide, Emu(nx), Emu(ny), Emu(node), Emu(node), fill=accent,
              shape=MSO_SHAPE.OVAL)
        _icon(slide, Emu(nx), Emu(ny), Emu(node), m.get("icon", T.ICON["check"]),
              color=T.BLANCO)

        # Bloque [ano][hito][texto] a la derecha del nodo (consistente).
        blk_w = int(Inches(4.4))
        bx = cx + node // 2 + int(Inches(0.65))
        align = PP_ALIGN.LEFT
        yx = bx
        sy2 = node_cy - int(Inches(0.95))
        # Etiqueta de ano como la eyebrow de la portada (pildora crema + mono azul).
        _pill(slide, Emu(yx), Emu(sy2), Inches(1.15), Inches(0.4),
              str(m.get("year", "")), fill=T.AMARILLO_TINT,
              text_color=T.AZUL_OSCURO, size=Pt(12))
        _text(slide, Emu(bx), Emu(sy2 + int(Inches(0.56))), Emu(blk_w), Inches(0.5),
              [[(m.get("label", ""), {"size": Pt(19), "color": T.AZUL_OSCURO,
                                      "font": T.FONT_HEAD})]], align=align)
        if m.get("text"):
            _text(slide, Emu(bx), Emu(sy2 + int(Inches(1.12))), Emu(blk_w),
                  Inches(1.2),
                  [[(m["text"], {"size": Pt(13), "color": T.GRIS_SUAVE,
                                 "font": T.FONT_BODY})]], align=align,
                  line_spacing=1.3)
        created.append(slide)
    return created


def add_comparison(prs, title, columns, rows, subtitle="", page=None, section="",
                   highlight=1):
    """Tabla comparativa con checks. columns: [etiqueta_col0, col1, col2, ...];
    rows: [{"label": str, "marks": [bool, ...]}]. `highlight` = indice de columna
    de datos a resaltar (1 = primera empresa). subtitle: serif opcional."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.4))
    content_top = tb
    if subtitle:
        sy = Emu(int(tb) + int(Inches(0.12)))
        _text(slide, MARGIN, sy, Inches(9.0), Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        content_top = Emu(int(sy) + int(Inches(0.5)))

    x0 = int(MARGIN)
    top = max(int(Inches(2.75)), int(content_top) + int(Inches(0.3)))
    bottom_limit = int(T.SLIDE_H) - int(Inches(0.85))
    label_w = int(Inches(3.0))
    ncol = max(len(columns) - 1, 1)
    col_w = (int(CONTENT_W) - label_w) // ncol
    head_h = int(Inches(0.64))
    nrows = max(len(rows), 1)
    row_h = (bottom_limit - top - head_h) // nrows   # ajustado para dejar margen
    table_h = head_h + nrows * row_h
    hl_x = x0 + label_w + (highlight - 1) * col_w

    def col_center(ci):
        return x0 + label_w + ci * col_w + col_w // 2

    # Columna propia como bloque azul marino (destacada). Sin barra de cabecera:
    # los titulos van en negrita sobre fondo transparente.
    _rect(slide, Emu(hl_x), Emu(top), Emu(col_w), Emu(table_h), fill=T.AZUL_OSCURO,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    _text(slide, Emu(x0 + int(Inches(0.3))), Emu(top),
          Emu(label_w - int(Inches(0.3))), Emu(head_h),
          [[(columns[0], {"size": Pt(13.5), "bold": True, "color": T.AZUL_OSCURO,
                          "font": T.FONT_TITLE})]], anchor=MSO_ANCHOR.MIDDLE)
    for ci, name in enumerate(columns[1:]):
        our = (ci == (highlight - 1))
        _text(slide, Emu(x0 + label_w + ci * col_w), Emu(top), Emu(col_w),
              Emu(head_h), [[(name, {"size": Pt(13), "bold": True,
                             "color": (T.AMARILLO if our else T.AZUL_OSCURO),
                             "font": T.FONT_TITLE})]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Divisor bajo la cabecera: a todo el ancho, azul marino, mismo grosor.
    hy = top + head_h
    _rect(slide, Emu(x0), Emu(hy), CONTENT_W, Pt(1), fill=T.AZUL_OSCURO)
    # Filas.
    for ri, row in enumerate(rows):
        ry = top + head_h + ri * row_h
        _text(slide, Emu(x0 + int(Inches(0.3))), Emu(ry),
              Emu(label_w - int(Inches(0.3))), Emu(row_h),
              [[(row["label"], {"size": Pt(13), "color": T.GRIS_TEXTO,
                                "font": T.FONT_BODY})]], anchor=MSO_ANCHOR.MIDDLE)
        # Divisores a todo el ancho, azul marino.
        if ri < len(rows) - 1:
            _rect(slide, Emu(x0), Emu(ry + row_h), CONTENT_W, Pt(1),
                  fill=T.AZUL_OSCURO)
        for ci, ok in enumerate(row.get("marks", [])):
            if not ok:
                continue
            our = (ci == (highlight - 1))
            sz = int(Inches(0.56))
            _icon(slide, Emu(col_center(ci) - sz // 2),
                  Emu(ry + (row_h - sz) // 2), Emu(sz), T.ICON["tick"],
                  color=(T.AMARILLO if our else T.AZUL))
    _pagenum(slide, page)
    return slide


def _letter_badge(slide, x, y, size, letter, color=None):
    """Circulo con inicial (A, B, C...) al estilo de las galerias-lista."""
    color = color if color is not None else T.AZUL_OSCURO
    _rect(slide, x, y, size, size, line=color, line_w=Pt(1.5), shape=MSO_SHAPE.OVAL)
    _text(slide, x, y, size, size,
          [[(letter, {"size": Pt(16), "bold": True, "color": color,
                      "font": T.FONT_HEAD})]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _num_badge(slide, x, y, size, num):
    """Circulo RELLENO navy con un ordinal dorado (filas de add_pricing).

    Distinto de _letter_badge, que es de contorno y usa el color del trazo para
    el texto.
    """
    _rect(slide, x, y, size, size, fill=T.AZUL_OSCURO, shape=MSO_SHAPE.OVAL)
    _text(slide, x, y, size, size,
          [[(num, {"size": Pt(12), "bold": True, "color": T.AMARILLO,
                   "font": T.FONT_NUM})]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _image_grid(slide, images, x, y, w, h, cols, rows, gap=Inches(0.1), radius=0.0):
    """Rejilla de imagenes (cover) en una region dada."""
    cw = Emu((int(w) - (cols - 1) * int(gap)) // cols)
    ch = Emu((int(h) - (rows - 1) * int(gap)) // rows)
    for i, img in enumerate(images[:cols * rows]):
        r, c = divmod(i, cols)
        ix = Emu(int(x) + c * (int(cw) + int(gap)))
        iy = Emu(int(y) + r * (int(ch) + int(gap)))
        _img(slide, img, ix, iy, cw, ch, radius=radius)


def add_gallery_list(prs, title, subtitle, items, images, page=None):
    """Galeria + lista: a la izquierda titulo, subtitulo y lista A/B/C/D; a la
    derecha rejilla de imagenes a sangre (2 o 4 segun cuantas se pasen)."""
    slide = _slide(prs)
    _topbar(slide)
    # Rejilla de imagenes a la derecha, a sangre.
    gx = Inches(6.95)
    gw = Emu(int(T.SLIDE_W) - int(gx))
    imgs = images[:4]
    if len(imgs) >= 4:
        _image_grid(slide, imgs, gx, 0, gw, T.SLIDE_H, 2, 2, gap=Inches(0.06))
    else:
        _image_grid(slide, imgs[:2], gx, 0, gw, T.SLIDE_H, 2, 1, gap=Inches(0.06))

    tw = Emu(int(gx) - int(MARGIN) - int(Inches(0.5)))
    tb = _stack_title(slide, _emph_runs(title, Pt(30)), MARGIN, Inches(1.4), tw, 30)
    y = Emu(int(tb) + int(Inches(0.12)))
    if subtitle:
        _text(slide, MARGIN, y, tw, Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True, "color": T.GRIS_SUAVE,
                            "font": T.FONT_TITLE_EMPH})]])
    items = items[:4]
    ny = Inches(3.05)
    step = Emu((int(Inches(6.7)) - int(ny)) // max(len(items), 1))
    badge = Inches(0.6)
    for i, it in enumerate(items):
        head, txt = (it if isinstance(it, (list, tuple)) else (it, ""))
        iy = Emu(int(ny) + i * int(step))
        _letter_badge(slide, MARGIN, iy, badge, chr(65 + i))
        hx = Emu(int(MARGIN) + int(badge) + int(Inches(0.3)))
        hw = Emu(int(tw) - int(badge) - int(Inches(0.3)))
        # Titulo + descripcion como pareja compacta (gap minimo), alineados con
        # el badge.
        runs = [[(head, {"size": Pt(15), "color": T.AZUL_OSCURO,
                         "font": T.FONT_HEAD})]]
        if txt:
            runs.append([(txt, {"size": Pt(11.5), "color": T.GRIS_SUAVE,
                                "font": T.FONT_BODY})])
        _text(slide, hx, iy, hw, badge, runs, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.12, space_after=Pt(3))
    _pagenum(slide, page)
    return slide


def add_gallery_text(prs, title, heading, body, images, subtitle="", page=None,
                     section=""):
    """Texto + galeria: parrafo a la izquierda y 3 imagenes verticales a la
    derecha (enmarcadas, esquinas sutiles). subtitle: serif opcional."""
    slide = _slide(prs)
    _topbar(slide, section)
    tw = Inches(4.9)
    tb = _stack_title(slide, _emph_runs(title, Pt(28)), MARGIN, Inches(1.4), tw, 28)
    y = Emu(int(tb) + int(Inches(0.15)))
    if subtitle:
        _text(slide, MARGIN, y, tw, Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        y = Emu(int(y) + int(Inches(0.6)))
    if heading:
        _text(slide, MARGIN, y, tw, Inches(0.45),
              [[(heading, {"size": Pt(16), "color": T.AZUL_OSCURO,
                           "font": T.FONT_HEAD})]])
        y = Emu(int(y) + int(Inches(0.55)))
    if isinstance(body, str):
        body = [body]
    paras = [[(b, {"size": Pt(13), "color": T.GRIS_SUAVE, "font": T.FONT_BODY})]
             for b in body]
    _text(slide, MARGIN, y, tw, Inches(3.4), paras, line_spacing=1.35,
          space_after=Pt(6))
    gx = Inches(6.3)
    gw = Emu(int(T.SLIDE_W) - int(gx) - int(MARGIN))
    _image_grid(slide, images[:3], gx, Inches(1.55), gw, Inches(5.3), 3, 1,
                gap=Inches(0.06), radius=0.025)
    _pagenum(slide, page)
    return slide


def add_gallery_mosaic(prs, title, images, subtitle="", page=None, section=""):
    """Mosaico de 5 imagenes: 2x2 pequenas a la izquierda + 1 grande a la
    derecha. Titulo arriba."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _stack_title(slide, _emph_runs(title, Pt(28)), MARGIN, Inches(1.4),
                      CONTENT_W, 28)
    gy = Inches(2.55)
    if subtitle:
        _text(slide, MARGIN, Emu(int(tb) + int(Inches(0.12))), Inches(9.0),
              Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        gy = Inches(3.05)
    gh = Emu(int(Inches(6.9)) - int(gy))
    gap = Inches(0.06)
    left_w = Inches(5.0)
    imgs = (images + [images[-1]] * 5)[:5] if images else []
    # 2x2 a la izquierda.
    _image_grid(slide, imgs[:4], MARGIN, gy, left_w, gh, 2, 2, gap=gap,
                radius=0.02)
    # 1 grande a la derecha.
    bx = Emu(int(MARGIN) + int(left_w) + int(gap))
    bw = Emu(int(T.SLIDE_W) - int(MARGIN) - int(bx))
    _img(slide, imgs[4], bx, gy, bw, gh, radius=0.02)
    _pagenum(slide, page)
    return slide


def _tint(slide, x, y, w, h, color=None, alpha=55000, radius=None):
    """Rect de color semitransparente (velo) sobre una foto."""
    color = color if color is not None else T.AZUL_OSCURO
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = _rect(slide, x, y, w, h, shape=shape, radius=radius)
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    fill = parse_xml('<a:solidFill xmlns:a="%s"><a:srgbClr val="%s">'
                     '<a:alpha val="%d"/></a:srgbClr></a:solidFill>'
                     % (_A, str(color), alpha))
    ln = spPr.find(qn("a:ln"))
    ln.addprevious(fill) if ln is not None else spPr.append(fill)
    return sp


def _glass_card(slide, image_path, x, y, w, h, radius=0.05, blur=55000):
    """Coloca dentro del rectangulo (x,y,w,h) una copia de la foto de fondo
    recortada EXACTAMENTE a esa zona (mismo encuadre 'cover' que el fondo a
    sangre) y desenfocada (`a:blur`). Simula el backdrop-blur del glassmorphism:
    a traves del cristal se ve el fondo borroso. Devuelve la imagen."""
    W, H = int(T.SLIDE_W), int(T.SLIDE_H)
    iw, ih = Image.open(image_path).size
    target, ratio = W / H, iw / ih
    # Ventana de imagen visible tras el recorte 'cover' del fondo (igual _img).
    if ratio > target:
        cxc, cyc = (1 - target / ratio) / 2, 0.0
    else:
        cxc, cyc = 0.0, (1 - ratio / target) / 2
    vis_l, vis_r, vis_t, vis_b = cxc, 1 - cxc, cyc, 1 - cyc
    fx0, fx1 = int(x) / W, (int(x) + int(w)) / W
    fy0, fy1 = int(y) / H, (int(y) + int(h)) / H
    img_l = vis_l + fx0 * (vis_r - vis_l)
    img_r = vis_l + fx1 * (vis_r - vis_l)
    img_t = vis_t + fy0 * (vis_b - vis_t)
    img_b = vis_t + fy1 * (vis_b - vis_t)
    pic = slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
    pic.crop_left, pic.crop_right = img_l, 1 - img_r
    pic.crop_top, pic.crop_bottom = img_t, 1 - img_b
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    val = int(radius * 100000)
    spPr.append(parse_xml(
        '<a:prstGeom xmlns:a="%s" prst="roundRect"><a:avLst>'
        '<a:gd name="adj" fmla="val %d"/></a:avLst></a:prstGeom>' % (_A, val)))
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="%s"><a:blur rad="%d"/></a:effectLst>'
        % (_A, int(blur))))
    return pic


def add_feature_band(prs, title, features, image, intro="", subtitle="",
                     page=None, section=""):
    """Banda de imagen (con velo azul marino) y 4 columnas de features: icono
    dorado + titulo + texto. features: [{icon, head, text}]. intro/subtitle
    opcionales."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.35))
    content_top = tb
    if subtitle:
        sy = Emu(int(tb) + int(Inches(0.12)))
        _text(slide, MARGIN, sy, Inches(9.0), Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        content_top = Emu(int(sy) + int(Inches(0.5)))
    band_y = Emu(max(int(Inches(2.5)), int(content_top) + int(Inches(0.25))))
    band_h = Inches(2.95)
    _img(slide, image, 0, band_y, T.SLIDE_W, band_h, radius=0)
    _tint(slide, 0, band_y, T.SLIDE_W, band_h, color=T.AZUL_OSCURO, alpha=74000)
    feats = features[:4]
    m = max(len(feats), 1)
    col_w = Emu((int(T.SLIDE_W) - 2 * int(MARGIN)) // m)
    for i, f in enumerate(feats):
        cx = Emu(int(MARGIN) + i * int(col_w))
        if i > 0:  # divisor sutil dorado
            _tint(slide, cx, Emu(int(band_y) + int(Inches(0.55))), Pt(1.4),
                  Emu(int(band_h) - int(Inches(1.1))), color=T.AMARILLO, alpha=55000)
        isz = Inches(0.6)
        _icon(slide, Emu(int(cx) + (int(col_w) - int(isz)) // 2),
              Emu(int(band_y) + int(Inches(0.58))), isz,
              f.get("icon", T.ICON["check"]), color=T.AMARILLO)
        _text(slide, Emu(int(cx) + int(Inches(0.2))),
              Emu(int(band_y) + int(Inches(1.22))),
              Emu(int(col_w) - int(Inches(0.4))), Inches(0.45),
              [[(f.get("head", ""), {"size": Pt(15), "color": T.BLANCO,
                                     "font": T.FONT_HEAD})]], align=PP_ALIGN.CENTER)
        if f.get("text"):
            _text(slide, Emu(int(cx) + int(Inches(0.32))),
                  Emu(int(band_y) + int(Inches(1.62))),
                  Emu(int(col_w) - int(Inches(0.64))), Inches(1.0),
                  [[(f["text"], {"size": Pt(11), "color": T.BLANCO,
                                 "font": T.FONT_BODY})]], align=PP_ALIGN.CENTER,
                  line_spacing=1.22)
    if intro:
        _text(slide, MARGIN, Emu(int(band_y) + int(band_h) + int(Inches(0.32))),
              CONTENT_W, Inches(0.8),
              [[(intro, {"size": Pt(13), "color": T.GRIS_SUAVE,
                         "font": T.FONT_BODY})]], align=PP_ALIGN.CENTER,
              line_spacing=1.3)
    _pagenum(slide, page)
    return slide


def add_mission(prs, title, features, image, caption_title, caption_text,
                page=None, section=""):
    """Lista de features a la izquierda (chip + titulo + texto) e imagen a la
    derecha con pie en caja azul marino. features: [{icon, head, text}]."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.4))

    fx = MARGIN
    fw = Inches(5.3)
    fy0 = Emu(max(int(Inches(2.75)), int(tb) + int(GAP_AFTER_TITLE)))
    feats = features[:4]
    step = Emu((int(Inches(6.8)) - int(fy0)) // max(len(feats), 1))
    chip = Inches(0.64)
    for i, f in enumerate(feats):
        iy = Emu(int(fy0) + i * int(step))
        _icon_chip(slide, fx, iy, chip, f.get("icon", T.ICON["check"]))
        hx = Emu(int(fx) + int(Inches(0.95)))
        hw = Emu(int(fw) - int(Inches(0.95)))
        # Titulo + texto como pareja compacta centrada con el chip (igual que las
        # vinetas del slide 5).
        runs = [[(f.get("head", ""), {"size": Pt(14.5), "color": T.AZUL_OSCURO,
                                      "font": T.FONT_HEAD})]]
        if f.get("text"):
            runs.append([(f["text"], {"size": Pt(11.5), "color": T.GRIS_SUAVE,
                                      "font": T.FONT_BODY})])
        _text(slide, hx, iy, hw, chip, runs, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.12, space_after=Pt(3))

    rx = Inches(6.85)
    rw = Emu(int(T.SLIDE_W) - int(rx))
    ry = Inches(1.15)
    img_h = Inches(3.7)
    _img(slide, image, rx, ry, rw, img_h, radius=0)
    box_y = Emu(int(ry) + int(img_h))
    box_h = Emu(int(T.SLIDE_H) - int(box_y))
    _rect(slide, rx, box_y, rw, box_h, fill=T.AZUL_OSCURO)
    cpad = Inches(0.5)
    _stack_title(slide, _emph_runs(caption_title, Pt(26), color=T.BLANCO,
                                   emph_color=T.AMARILLO),
                 Emu(int(rx) + int(cpad)), Emu(int(box_y) + int(cpad)),
                 Emu(int(rw) - 2 * int(cpad)), 26, line_h=1.1)
    if caption_text:
        _text(slide, Emu(int(rx) + int(cpad)),
              Emu(int(box_y) + int(cpad) + int(Inches(0.7))),
              Emu(int(rw) - 2 * int(cpad)), Inches(1.2),
              [[(caption_text, {"size": Pt(12), "color": T.BLANCO,
                                "font": T.FONT_BODY})]], line_spacing=1.3)
    _pagenum(slide, page)
    return slide


def add_numbered_grid(prs, title, items, highlight="", page=None, section=""):
    """Rejilla de items numerados (numero mono + titulo + texto) en 3 columnas +
    una declaracion resaltada (marcador dorado). items: [(titulo, texto)] hasta 4."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.4))
    top = Emu(max(int(Inches(2.7)), int(tb) + int(GAP_AFTER_TITLE)))
    gap = Inches(0.5)
    col_w = Emu((int(CONTENT_W) - 2 * int(gap)) // 3)
    row_h = Inches(2.0)
    its = items[:4]
    for i, it in enumerate(its):
        h, t = (it if isinstance(it, (list, tuple)) else (it, ""))
        r, c = divmod(i, 3)
        x = Emu(int(MARGIN) + c * (int(col_w) + int(gap)))
        y = Emu(int(top) + r * (int(row_h) + int(Inches(0.3))))
        _text(slide, x, y, col_w, Inches(0.6),
              [[("%02d" % (i + 1), {"size": Pt(30), "bold": True,
                 "color": T.AZUL_OSCURO, "font": T.FONT_NUM})]])
        _text(slide, x, Emu(int(y) + int(Inches(0.66))), col_w, Inches(0.4),
              [[(h, {"size": Pt(15), "color": T.AZUL_OSCURO,
                     "font": T.FONT_HEAD})]])
        if t:
            _text(slide, x, Emu(int(y) + int(Inches(1.1))), col_w, Inches(0.8),
                  [[(t, {"size": Pt(11.5), "color": T.GRIS_SUAVE,
                         "font": T.FONT_BODY})]], line_spacing=1.25)
    if highlight:
        hx = Emu(int(MARGIN) + int(col_w) + int(gap))
        hy = Emu(int(top) + int(row_h) + int(Inches(0.3)))
        hw = Emu(int(MARGIN) + int(CONTENT_W) - int(hx))
        _rect(slide, hx, hy, hw, Inches(1.35), fill=T.AMARILLO,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        _text(slide, Emu(int(hx) + int(Inches(0.4))), hy,
              Emu(int(hw) - int(Inches(0.8))), Inches(1.35),
              [_emph_runs(highlight, Pt(18), color=T.AZUL_OSCURO)],
              anchor=MSO_ANCHOR.MIDDLE, line_h=1.18)
    _pagenum(slide, page)
    return slide


def add_thanks(prs, headline, contact, image, page=None):
    """Cierre 'gracias' sobre foto a sangre con panel claro: titular grande +
    datos de contacto (icono + etiqueta + valor). contact: [(icon, label, value)]."""
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
    px, py = Inches(1.0), Inches(1.35)
    pw = Emu(int(T.SLIDE_W) - 2 * int(px))
    ph = Emu(int(T.SLIDE_H) - int(py) - int(Inches(1.05)))
    card = _rect(slide, px, py, pw, ph, fill=T.BG,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    _soft_shadow(card, alpha=26000)
    pad = Inches(0.7)
    tx = Emu(int(px) + int(pad))
    _stack_title(slide, _emph_runs(headline, Pt(40)), tx,
                 Emu(int(py) + int(pad) + int(Inches(0.3))), Inches(4.8), 40,
                 line_h=1.08)
    # Columna de contacto a la derecha del panel.
    cx = Emu(int(px) + int(pw) // 2 + int(Inches(0.2)))
    cw = Emu(int(px) + int(pw) - int(pad) - int(cx))
    cy = Emu(int(py) + int(pad) + int(Inches(0.35)))
    cstep = Emu((int(ph) - 2 * int(pad)) // max(len(contact), 1))
    for icon, label, value in contact:
        _icon_chip(slide, cx, cy, Inches(0.5), icon)
        lx = Emu(int(cx) + int(Inches(0.72)))
        lw = Emu(int(cx) + int(cw) - int(lx))
        _text(slide, lx, cy, lw, Inches(0.28),
              [[(label.upper(), {"size": Pt(10), "color": T.AZUL_OSCURO,
                                 "font": T.FONT_MONO, "spacing": 120})]])
        _text(slide, lx, Emu(int(cy) + int(Inches(0.26))), lw, Inches(0.32),
              [[(value, {"size": Pt(13), "color": T.GRIS_TEXTO,
                         "font": T.FONT_HEAD})]])
        _rect(slide, lx, Emu(int(cy) + int(Inches(0.62))), lw, Pt(1),
              fill=T.GRIS_BORDE)
        cy = Emu(int(cy) + int(cstep))
    # Barra dorada al pie del panel.
    _rect(slide, px, Emu(int(py) + int(ph) - int(Inches(0.14))), pw, Inches(0.14),
          fill=T.AMARILLO)
    return slide


def add_hero_card(prs, title, body, cta, image, page=None):
    """Hero con foto a sangre y tarjeta clara (centrada a la izquierda): titular
    centrado, parrafo y boton al pie (mismo estilo de pill que el resto)."""
    slide = _slide(prs, bg=T.AZUL_OSCURO)
    _img(slide, image, 0, 0, T.SLIDE_W, T.SLIDE_H, radius=0)
    cw, chh = Inches(4.8), Inches(5.4)
    cxp = Inches(0.95)
    cyp = (int(T.SLIDE_H) - int(chh)) // 2
    card = _rect(slide, cxp, Emu(cyp), cw, chh, fill=T.BG,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    _soft_shadow(card, alpha=28000)
    pad = int(Inches(0.55))
    inx = Emu(int(cxp) + pad)
    inw = Emu(int(cw) - 2 * pad)
    # Titulo centrado arriba.
    tb = _stack_title(slide, _emph_runs(title, Pt(27)), inx,
                      Emu(cyp + int(Inches(0.85))), inw, 27, line_h=1.14,
                      align=PP_ALIGN.CENTER)
    # Parrafo centrado bajo el titulo.
    _text(slide, inx, Emu(int(tb) + int(Inches(0.4))), inw, Inches(1.9),
          [[(body, {"size": Pt(12), "color": T.GRIS_SUAVE, "font": T.FONT_BODY})]],
          align=PP_ALIGN.CENTER, line_spacing=1.45)
    # Boton al pie del card, con el estilo del resto (pill azul + mono).
    bw, bh = Inches(2.5), Inches(0.55)
    bx = Emu(int(cxp) + (int(cw) - int(bw)) // 2)
    by = Emu(cyp + int(chh) - pad - int(bh))
    _pill(slide, bx, by, bw, bh, cta, fill=T.AZUL, text_color=T.BLANCO,
          size=Pt(12))
    _white_pagenum(slide, page)
    return slide


def add_value_cards(prs, title, image, cards, subtitle="", page=None, section=""):
    """Imagen alta a la izquierda + titulo arriba-derecha + fila de tarjetas de
    acento ABAJO (alineadas con la base de la imagen), cada una con flecha en
    circulo. cards: {"head","text"} o {"value","label"} (tarjeta-cifra), hasta 3."""
    slide = _slide(prs)
    _topbar(slide, section)
    cards_bottom = int(T.SLIDE_H) - int(Inches(0.7))
    iy = int(Inches(1.5))
    iw = int(Inches(4.4))
    _img(slide, image, MARGIN, Emu(iy), Emu(iw), Emu(cards_bottom - iy),
         radius=0.03)
    rx = int(MARGIN) + iw + int(Inches(0.6))
    rw = int(T.SLIDE_W) - int(MARGIN) - rx
    tb = _stack_title(slide, _emph_runs(title, Pt(30)), Emu(rx), Inches(1.6),
                      Emu(rw), 30, line_h=1.12)
    if subtitle:
        _text(slide, Emu(rx), Emu(int(tb) + int(Inches(0.12))), Emu(rw),
              Inches(0.5), [[(subtitle, {"size": Pt(14), "italic": True,
              "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])

    cards = cards[:3]
    n = max(len(cards), 1)
    cgap = int(Inches(0.25))
    cw = (rw - (n - 1) * cgap) // n
    ch = int(Inches(2.7))
    cy = cards_bottom - ch
    pad = int(Inches(0.32))
    circ = int(Inches(0.56))

    def arrow_circle(ax, ay):
        _rect(slide, Emu(ax), Emu(ay), Emu(circ), Emu(circ), fill=T.AZUL_OSCURO,
              shape=MSO_SHAPE.OVAL)
        _icon(slide, Emu(ax), Emu(ay), Emu(circ), T.ICON["arrow"], color=T.BLANCO)

    for i, c in enumerate(cards):
        cxx = rx + i * (cw + cgap)
        stat = "value" in c
        _rect(slide, Emu(cxx), Emu(cy), Emu(cw), Emu(ch),
              fill=(T.AMARILLO_TINT if stat else T.AMARILLO),
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        txx = cxx + pad
        tww = cw - 2 * pad
        arrow_circle(cxx + cw - pad - circ, cy + pad)   # flecha arriba-derecha
        hy2 = cy + pad + circ + int(Inches(0.18))
        if stat:
            # Cifra grande destacada + etiqueta debajo.
            _text(slide, Emu(txx), Emu(hy2), Emu(tww), Inches(0.7),
                  [[(c["value"], {"size": Pt(32), "bold": True,
                     "color": T.AZUL_OSCURO, "font": T.FONT_NUM})]])
            _text(slide, Emu(txx), Emu(hy2 + int(Inches(0.78))), Emu(tww),
                  Inches(0.8),
                  [[(c.get("label", ""), {"size": Pt(11.5), "color": T.AZUL_OSCURO,
                                          "font": T.FONT_BODY})]], line_spacing=1.2)
        else:
            _text(slide, Emu(txx), Emu(hy2), Emu(tww), Inches(0.6),
                  [[(c.get("head", ""), {"size": Pt(15), "color": T.AZUL_OSCURO,
                                         "font": T.FONT_HEAD})]], line_spacing=1.1)
            _text(slide, Emu(txx), Emu(hy2 + int(Inches(0.54))), Emu(tww),
                  Inches(1.0),
                  [[(c.get("text", ""), {"size": Pt(11), "color": T.AZUL_OSCURO,
                                         "font": T.FONT_BODY})]], line_spacing=1.25)
    _pagenum(slide, page)
    return slide


def add_overview(prs, title, subhead, body, points, images, subtitle="",
                 page=None, section=""):
    """Resumen a 3 columnas: texto a la izquierda, 2 imagenes apiladas en medio,
    y 2 puntos numerados (01/ 02/) a la derecha. points: [(head, text)].
    subtitle: serif opcional."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.4))
    content_top = tb
    if subtitle:
        sy = Emu(int(tb) + int(Inches(0.12)))
        _text(slide, MARGIN, sy, Inches(9.0), Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE, "font": T.FONT_TITLE_EMPH})]])
        content_top = Emu(int(sy) + int(Inches(0.5)))
    top = Emu(max(int(Inches(2.7)), int(content_top) + int(Inches(0.3))))

    lw = Inches(3.15)
    _text(slide, MARGIN, top, lw, Inches(0.95),
          [[(subhead, {"size": Pt(15), "color": T.AZUL_OSCURO,
                       "font": T.FONT_HEAD})]], line_spacing=1.2)
    if isinstance(body, str):
        body = [body]
    paras = [[(b, {"size": Pt(11.5), "color": T.GRIS_SUAVE, "font": T.FONT_BODY})]
             for b in body]
    _text(slide, MARGIN, Emu(int(top) + int(Inches(1.05))), lw, Inches(3.4),
          paras, line_spacing=1.3, space_after=Pt(8))

    mx = Emu(int(MARGIN) + int(lw) + int(Inches(0.5)))
    mw = Inches(3.55)
    igap = int(Inches(0.06))   # mismo gap minimo que las galerias
    imh = Emu((int(Inches(4.0)) - igap) // 2)
    for i, img in enumerate(images[:2]):
        _img(slide, img, mx, Emu(int(top) + i * (int(imh) + igap)),
             mw, imh, radius=0.025)

    rx = Emu(int(mx) + int(mw) + int(Inches(0.55)))
    rw = Emu(int(T.SLIDE_W) - int(MARGIN) - int(rx))
    pts = points[:2]
    pstep = Emu(int(Inches(4.0)) // max(len(pts), 1))
    for i, pt in enumerate(pts):
        h, t = (pt if isinstance(pt, (list, tuple)) else (pt, ""))
        py = Emu(int(top) + i * int(pstep))
        _text(slide, rx, py, Inches(1.2), Inches(0.7),
              [[("%02d/" % (i + 1), {"size": Pt(34), "bold": True,
                 "color": T.AZUL, "font": T.FONT_NUM})]])
        hx = Emu(int(rx) + int(Inches(1.3)))
        hw = Emu(int(rx) + int(rw) - int(hx))
        # Head + texto como pareja compacta (gap minimo, alineada con el numero).
        runs = [[(h, {"size": Pt(15), "color": T.AZUL_OSCURO,
                      "font": T.FONT_HEAD})]]
        if t:
            runs.append([(t, {"size": Pt(11.5), "color": T.GRIS_SUAVE,
                              "font": T.FONT_BODY})])
        _text(slide, hx, py, hw, Inches(1.5), runs, line_spacing=1.14,
              space_after=Pt(3))
    _pagenum(slide, page)
    return slide


def add_product(prs, title, stat, kicker, rows, page=None, section=""):
    """Producto: titulo grande + cifra + kicker a la izquierda; a la derecha
    filas de [titulo+texto | imagen | parrafo]. rows: [{head, text, image, desc}]."""
    slide = _slide(prs)
    _topbar(slide, section)
    lw = Inches(3.0)
    _stack_title(slide, _emph_runs(title, Pt(34)), MARGIN, Inches(1.55), lw, 34,
                 line_h=1.05)
    _text(slide, MARGIN, Inches(3.35), lw, Inches(1.0),
          [[(stat, {"size": Pt(50), "bold": True, "color": T.AZUL_OSCURO,
                    "font": T.FONT_NUM})]])
    _text(slide, MARGIN, Inches(4.5), lw, Inches(0.4),
          [[("< " + kicker + " >", {"size": Pt(10.5), "color": T.GRIS_SUAVE,
             "font": T.FONT_MONO, "spacing": 60})]])

    rx = Emu(int(MARGIN) + int(lw) + int(Inches(0.5)))
    rw = int(T.SLIDE_W) - int(MARGIN) - int(rx)
    rows = rows[:3]
    top = Inches(1.6)
    rh = Emu(int(Inches(5.3)) // max(len(rows), 1))
    cA = int(rw * 0.34)
    cImg = int(rw * 0.30)
    for i, r in enumerate(rows):
        ry = Emu(int(top) + i * int(rh))
        _text(slide, rx, ry, Emu(cA - int(Inches(0.3))), Inches(0.5),
              [[(r.get("head", ""), {"size": Pt(17), "color": T.AZUL_OSCURO,
                                     "font": T.FONT_HEAD})]])
        if r.get("text"):
            _text(slide, rx, Emu(int(ry) + int(Inches(0.5))),
                  Emu(cA - int(Inches(0.3))), Inches(0.9),
                  [[(r["text"], {"size": Pt(10.5), "color": T.GRIS_SUAVE,
                                 "font": T.FONT_BODY})]], line_spacing=1.2)
        imx = Emu(int(rx) + cA)
        _img(slide, r["image"], imx, ry, Emu(cImg - int(Inches(0.3))),
             Emu(int(rh) - int(Inches(0.35))), radius=0.03)
        cxp = Emu(int(imx) + cImg)
        cwp = Emu(int(rx) + rw - int(cxp))
        _text(slide, cxp, ry, cwp, Emu(int(rh) - int(Inches(0.35))),
              [[(r.get("desc", ""), {"size": Pt(10.5), "color": T.GRIS_SUAVE,
                                     "font": T.FONT_BODY})]], line_spacing=1.25)
    _pagenum(slide, page)
    return slide


def _pricing_page(prs, title, rows, ordinal, subtitle, section, page,
                  texto_total, note):
    """Pinta UNA pagina del desglose. texto_total=None -> sin tarjeta (pagina
    no final): las filas ocupan todo el ancho."""
    slide = _slide(prs)
    _topbar(slide, section)
    tb = _title(slide, title, y=Inches(1.35))
    top = max(int(Inches(2.5)), int(tb) + int(GAP_AFTER_TITLE))
    if subtitle:
        _text(slide, MARGIN, Emu(int(tb) + int(Inches(0.1))), Inches(9.0),
              Inches(0.5),
              [[(subtitle, {"size": Pt(14), "italic": True,
                            "color": T.GRIS_SUAVE,
                            "font": T.FONT_TITLE_EMPH})]])
        top += int(Inches(0.45))

    bottom = int(Inches(6.3)) if note else int(Inches(6.6))
    block_h = bottom - top

    con_total = texto_total is not None
    gap_col = int(Inches(0.5))
    rows_w = int(0.60 * int(CONTENT_W)) if con_total else int(CONTENT_W)

    # Alto de fila topado: sin el tope, una sola partida ocuparia todo el hueco.
    # El bloque resultante se centra verticalmente entre `top` y `bottom`.
    k = len(rows)
    gap_r = int(Inches(0.18))
    row_h = max(int(Inches(0.4)),
                min(int(ROW_H_MAX), (block_h - (k - 1) * gap_r) // k))
    used_h = k * row_h + (k - 1) * gap_r
    rows_top = top + (block_h - used_h) // 2

    badge = int(Inches(0.34))
    pad = int(Inches(0.28))
    # 1.9in basta: el importe mas largo imaginable ("999.999,99 EUR" en Geist
    # Mono a 15pt) mide 1.75in. Lo que sobra es ancho para el concepto, que es
    # quien envuelve a dos lineas.
    amount_w = int(Inches(1.9))
    concept_x = int(MARGIN) + pad + badge + pad
    concept_w = int(MARGIN) + rows_w - amount_w - int(Inches(0.5)) - concept_x

    for i, (concepto, importe) in enumerate(rows):
        y = rows_top + i * (row_h + gap_r)
        card = _rect(slide, MARGIN, Emu(y), Emu(rows_w), Emu(row_h),
                     fill=T.BLANCO, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                     radius=0.06)
        _soft_shadow(card, alpha=10000)
        _num_badge(slide, Emu(int(MARGIN) + pad),
                   Emu(y + (row_h - badge) // 2), Emu(badge),
                   "%d" % (ordinal + i))
        _text(slide, Emu(concept_x), Emu(y), Emu(concept_w), Emu(row_h),
              [[(concepto, {"size": Pt(14.5), "color": T.AZUL_OSCURO,
                            "font": T.FONT_HEAD})]],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(int(MARGIN) + rows_w - amount_w - int(Inches(0.35))),
              Emu(y), Emu(amount_w), Emu(row_h),
              [[(_fmt_eur(importe), {"size": Pt(15), "bold": True,
                                     "color": T.AZUL_OSCURO,
                                     "font": T.FONT_NUM})]],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # La tarjeta se alinea con el BLOQUE REAL de filas (rows_top/used_h), no con
    # el hueco completo: si no, con pocas partidas quedaria descuadrada.
    if con_total:
        card_x = int(MARGIN) + rows_w + gap_col
        card_w = int(CONTENT_W) - rows_w - gap_col
        card_h = max(int(used_h), int(CARD_H_MIN))
        card_y = int(rows_top) + (int(used_h) - card_h) // 2
        card = _rect(slide, Emu(card_x), Emu(card_y), Emu(card_w),
                     Emu(card_h), fill=T.AMARILLO,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        _soft_shadow(card, alpha=9000)
        _text(slide, Emu(card_x), Emu(card_y + int(Inches(0.35))), Emu(card_w),
              Inches(0.35),
              [[("TOTAL ESTIMADO", {"size": Pt(11), "color": T.AZUL_OSCURO,
                                    "font": T.FONT_MONO, "spacing": 120})]],
              align=PP_ALIGN.CENTER)
        cifra_h = int(Inches(1.2))
        _text(slide, Emu(card_x), Emu(card_y + (card_h - cifra_h) // 2),
              Emu(card_w), Emu(cifra_h),
              [[(texto_total, {"size": Pt(38), "bold": True,
                               "color": T.AZUL_OSCURO, "font": T.FONT_NUM})]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Emu(card_x), Emu(card_y + card_h - int(Inches(0.7))),
              Emu(card_w), Inches(0.35),
              [[("IVA no incluido", {"size": Pt(12), "italic": True,
                                     "color": T.AZUL_OSCURO,
                                     "font": T.FONT_TITLE_EMPH})]],
              align=PP_ALIGN.CENTER)

    if note:
        _text(slide, MARGIN, Emu(bottom + int(Inches(0.25))),
              Emu(int(CONTENT_W) - int(Inches(1.4))),
              Inches(0.6),
              [[(note, {"size": Pt(10), "italic": True, "color": T.GRIS_SUAVE,
                        "font": T.FONT_BODY})]], line_spacing=1.25)

    _pagenum(slide, page)
    return slide


def add_pricing(prs, title, rows, note="", total=None, subtitle="", page=None,
                section=""):
    """Desglose de presupuesto: partidas numeradas + tarjeta de total + nota.

    rows  : [(concepto, importe)] con importe int/float. Maximo 10 partidas.
    total : None -> se suma a partir de rows. Un str se pinta tal cual
            ("A convenir"). NO acepta un numero: el total nunca debe poder
            contradecir al desglose.
    page  : el CONTADOR, no su valor. page=n, no page=n(). Se invoca una vez
            por pagina generada.

    MULTIPAGINA: con 6-10 partidas genera dos diapositivas (la primera a ancho
    completo, sin tarjeta) y el total va en la ultima. Devuelve SIEMPRE la lista
    de slides, tenga una o dos.
    """
    if not rows:
        raise ValueError("add_pricing: 'rows' no puede estar vacia")
    if len(rows) > MAX_PRICING_ROWS:
        raise ValueError(
            "add_pricing: %d partidas; el maximo es %d. Agrupa partidas o usa "
            "dos diapositivas: no se truncan en silencio."
            % (len(rows), MAX_PRICING_ROWS))
    if total is not None and not isinstance(total, str):
        raise TypeError(
            "add_pricing: 'total' debe ser None o str (p.ej. 'A convenir'); "
            "los importes se suman a partir de 'rows'")
    if page is not None and not callable(page):
        raise TypeError(
            "add_pricing: 'page' debe ser el contador (page=n), no su valor "
            "(page=n())")

    rows = [tuple(r) for r in rows]
    if total is not None:
        texto_total = total
    else:
        # Sumamos centimos ya redondeados: asi el total SIEMPRE coincide con la
        # suma de los importes que se pintan en las filas.
        texto_total = _fmt_eur_centimos(sum(_centimos(r[1]) for r in rows))
    paginas = _split_rows(rows)
    slides = []
    ordinal = 1
    for i, chunk in enumerate(paginas):
        ultima = (i == len(paginas) - 1)
        slides.append(_pricing_page(
            prs, title, chunk, ordinal=ordinal, subtitle=subtitle,
            section=section, page=(page() if page is not None else None),
            texto_total=(texto_total if ultima else None),
            note=(note if ultima else "")))
        ordinal += len(chunk)
    return slides
